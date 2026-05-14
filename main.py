
from fastapi import FastAPI, Request, UploadFile, File, HTTPException
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pathlib import Path
import json, math, io, zipfile, csv, tempfile, shutil, os, hashlib, secrets, base64, hmac, time
from datetime import datetime
import matplotlib
import psycopg2

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).parent
APP_DIR = ROOT
STATIC = ROOT / "static"
UPLOADS = ROOT / "uploads"
GENERATED = Path("/tmp/generated_ppt")
DATA_DIR = Path("/tmp")
PROJECTS = DATA_DIR / "user_projects"
USERS = DATA_DIR / "users.json"
for p in (UPLOADS, GENERATED, PROJECTS):
    p.mkdir(exist_ok=True)

app = FastAPI(title="DyeFlow RS Web v42 Overflow Wave Engine")
app.mount("/static", StaticFiles(directory=str(STATIC)), name="static")

SECRET = os.environ.get("DYEFLOW_SECRET_KEY", "dyeflow-local-dev-secret-v28")

@app.get("/")
def index():
    return FileResponse(ROOT / "index.html")

@app.get("/api/health")
def health():
    return {"ok": True, "version": "v54_enterprise_report"}

# ---------- Auth helpers ----------
DATABASE_URL = os.getenv("DATABASE_URL")

def get_db_connection():
    return psycopg2.connect(DATABASE_URL)

def _load_users():
    conn = get_db_connection()
    cur = conn.cursor()

    cur.execute("""
        CREATE TABLE IF NOT EXISTS users (
            username TEXT PRIMARY KEY,
            password TEXT,
            name TEXT,
            email TEXT,
            role TEXT,
            active BOOLEAN,
            can_save BOOLEAN
        )
    """)

    conn.commit()

    cur.execute("""
        SELECT username, password, name, email, role, active, can_save
        FROM users
    """)

    rows = cur.fetchall()

    users = []
    for row in rows:
        users.append({
            "username": row[0],
            "password": row[1],
            "name": row[2],
            "email": row[3],
            "role": row[4],
            "active": row[5],
            "can_save": row[6]
        })

    cur.close()
    conn.close()

    return {"users": users}

def _save_users(data):
    conn = get_db_connection()
    cur = conn.cursor()


    for u in data["users"]:
        cur.execute("""
    INSERT INTO users
    (username, password, name, email, role, active, can_save)
    VALUES (%s,%s,%s,%s,%s,%s,%s)

    ON CONFLICT (username)
    DO UPDATE SET
        password = EXCLUDED.password,
        name = EXCLUDED.name,
        email = EXCLUDED.email,
        role = EXCLUDED.role,
        active = EXCLUDED.active,
        can_save = EXCLUDED.can_save
""", (
            u.get("username"),
            u.get("password") or u.get("password_hash"),
            u.get("name"),
            u.get("email"),
            u.get("role"),
            u.get("active"),
            u.get("can_save")
        ))

    conn.commit()
    cur.close()
    conn.close()
def _hash_password(password, salt=None):
    salt = salt or secrets.token_hex(16)
    dk = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 120000)
    return salt + "$" + dk.hex()

def _verify_password(password, stored):
    try:
        salt, hv = stored.split("$", 1)
        return hmac.compare_digest(_hash_password(password, salt), stored)
    except Exception:
        return False

def _token(username):
    payload = {"u": username, "t": int(time.time())}
    raw = base64.urlsafe_b64encode(json.dumps(payload).encode()).decode().rstrip("=")
    sig = hmac.new(SECRET.encode(), raw.encode(), hashlib.sha256).hexdigest()
    return raw + "." + sig

def _user_from_request(req: Request):
    auth = req.headers.get("authorization") or ""
    if not auth.lower().startswith("bearer "):
        return None
    token = auth.split(" ", 1)[1]
    try:
        raw, sig = token.split(".", 1)
        if not hmac.compare_digest(hmac.new(SECRET.encode(), raw.encode(), hashlib.sha256).hexdigest(), sig):
            return None
        padded = raw + "=" * (-len(raw) % 4)
        username = json.loads(base64.urlsafe_b64decode(padded.encode()).decode()).get("u")
        users = _load_users()["users"]
        for u in users:
            if u.get("username") == username and u.get("is_active", True):
                return {k:v for k,v in u.items() if k != "password_hash"}
    except Exception:
        return None
    return None

@app.post("/api/auth/register")
async def register(req: Request):
    data = await req.json()
    users = _load_users()
    if users["users"]:
        raise HTTPException(403, "Register is only available for the first admin user.")
    username = (data.get("username") or data.get("email") or "").strip()
    password = data.get("password") or ""
    if len(username) < 2 or len(password) < 6:
        raise HTTPException(400, "Username and minimum 6-character password required.")
    user = {
        "username": username,
        "email": data.get("email",""),
        "name": data.get("name",""),
        "role": "admin",
        "can_save": True,
        "is_active": True,
        "password_hash": _hash_password(password),
        "created_at": datetime.now().isoformat(timespec="seconds")
    }
    users["users"].append(user)
    _save_users(users)
    return {"token": _token(username), "user": {k:v for k,v in user.items() if k!="password_hash"}}

@app.post("/api/auth/login")
async def login(req: Request):
    data = await req.json()
    login = (data.get("login") or "").strip()
    password = data.get("password") or ""
    for u in _load_users().get("users", []):
        if (u.get("username") == login or u.get("email") == login) and u.get("is_active", True):
            if _verify_password(password, u.get("password_hash") or u.get("password","")):
                return {"token": _token(u["username"]), "user": {k:v for k,v in u.items() if k!="password_hash"}}
    raise HTTPException(401, "Invalid login.")

@app.get("/api/auth/me")
def me(req: Request):
    u = _user_from_request(req)
    return {"authenticated": bool(u), "user": u}

@app.get("/api/admin/users")
def admin_users(req: Request):
    u = _user_from_request(req)
    if not u or u.get("role") != "admin":
        raise HTTPException(403, "Admin only.")
    users = [{k:v for k,v in x.items() if k!="password_hash"} for x in _load_users().get("users", [])]
    return {"users": users}

@app.post("/api/admin/users")
async def admin_create_user(req: Request):
    u = _user_from_request(req)
    if not u or u.get("role") != "admin":
        raise HTTPException(403, "Admin only.")
    data = await req.json()
    users = _load_users()
    username = (data.get("username") or "").strip()
    if not username or not data.get("password"):
        raise HTTPException(400, "Username/password required.")
    if any(x.get("username") == username for x in users["users"]):
        raise HTTPException(400, "User already exists.")
    nu = {
        "username": username, "email": data.get("email",""), "name": data.get("name",""),
        "role": data.get("role","user"), "can_save": bool(data.get("can_save", True)),
        "is_active": bool(data.get("is_active", True)), "password_hash": _hash_password(data["password"]),
        "created_at": datetime.now().isoformat(timespec="seconds")
    }
    users["users"].append(nu); _save_users(users)
    return {"user": {k:v for k,v in nu.items() if k!="password_hash"}}

@app.patch("/api/admin/users/{username}")
async def admin_update_user(username: str, req: Request):
    u = _user_from_request(req)
    if not u or u.get("role") != "admin":
        raise HTTPException(403, "Admin only.")
    data = await req.json()
    users = _load_users()
    for x in users["users"]:
        if x.get("username") == username:
            for k in ("name","role","is_active","can_save"):
                if k in data: x[k] = data[k]
            if data.get("password"):
                x["password_hash"] = _hash_password(data["password"])
            _save_users(users)
            return {"user": {k:v for k,v in x.items() if k!="password_hash"}}
    raise HTTPException(404, "User not found.")

# ---------- Calculation / graph ----------
def N(v, d=0.0):
    try:
        if v is None or v == "": return float(d)
        return float(v)
    except Exception:
        return float(d)

def calc(project):
    steps = project.get("steps") or []
    machine = project.get("machine") or {}
    utilities = project.get("utilities") or {}
    fabric = max(N(project.get("fabric_kg"), 1), 0.001)
    flote = N(project.get("flote"), 0)
    # Water is calculated step-by-step so overflow can add extra water.
    water_m3 = 0.0
    total_power = N(machine.get("circulation_pump_power"))*N(machine.get("pump_ratio"),1) + N(machine.get("number_of_reel"))*N(machine.get("reel_power")) + N(machine.get("fan_power"))*N(machine.get("fan_ratio"),1)
    drain_time = N(machine.get("drain_time_min"),5)
    x=[]; y=[]; events=[]; chemical_rows=[]; legend=[]
    t=0.0; cur=None; letter_i=0
    letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    # base_process_water_l: fresh water taken into the process, after carry-over logic.
    # overflow_extra_water_l: additional overflow water used only for water/wastewater cost.
    # User rules:
    # - Fabric Wet: every step fresh water = (step flote - carry over) × fabric kg.
    # - Fabric Dry: first step fresh water = step flote × fabric kg; next steps = (step flote - carry over) × fabric kg.
    # - Overflow water must NOT affect heating or chemical consumption.
    base_process_water_l = 0.0
    overflow_extra_water_l = 0.0
    carry_over_l_saved = 0.0
    wastewater_l = 0.0
    active_electric_time_min = 0.0
    current_bath_l = 0.0
    # Water actually taken since the last drain. Wastewater can never exceed this amount.
    pending_wastewater_l = 0.0
    previous_step_drained = True
    # Heating energy uses nominal bath water only: flote × fabric kg per step.
    # Carry over and overflow water are intentionally excluded from heating.
    heating_kcal_raw = 0.0

    carry_over = max(N(project.get("carry_over"), 0), 0)
    fabric_status = str(project.get("fabric_status") or "Dry").strip().lower()

    def step_flote_ratio(step):
        sr = N(step.get("flote_ratio"), 0)
        return sr if sr > 0 else flote

    def step_nominal_bath_l(step):
        # Prefer the visible step water amount if it exists; otherwise calculate from flote ratio.
        amount = N(step.get("amount_of_flote"), 0)
        return amount if amount > 0 else max(step_flote_ratio(step), 0) * fabric

    def step_fresh_water_l(step, step_index):
        fr = max(step_flote_ratio(step), 0)
        if fabric_status == "wet" or step_index > 0:
            effective_fr = max(fr - carry_over, 0)
        else:
            effective_fr = fr
        fresh = effective_fr * fabric
        nominal = fr * fabric
        return fresh, max(nominal - fresh, 0)

    def step_takes_fresh_water(step_index):
        # First bath always takes water. After that, a new bath is taken only if the previous step drained.
        return step_index == 0 or previous_step_drained

    def overflow_extra_liters(step, overflow):
        fill = N(step.get("filling_time"), 0)
        bath_l = step_nominal_bath_l(step)
        # Existing overflow rule stays: additional water = (overflow time / filling time) × nominal step bath water.
        return (overflow / fill) * bath_l if overflow > 0 and fill > 0 and bath_l > 0 else 0.0

    def add_point(tt, yy):
        x.append(round(tt,3)); y.append(round(yy,3))

    for si, s in enumerate(steps):
        begin=N(s.get("beginning_temp"), 25)
        fill=N(s.get("filling_time"),0)
        cur = begin if cur is None else begin
        add_point(t, cur)
        if fill>0:
            t += fill
            add_point(t, cur)
        fill_end=t
        if step_takes_fresh_water(si):
            fresh_water_l, saved_l = step_fresh_water_l(s, si)
            current_bath_l = step_nominal_bath_l(s)
        else:
            # If previous step was not drained, this step continues with the same bath: no new water intake.
            fresh_water_l, saved_l = 0.0, 0.0
        base_process_water_l += fresh_water_l
        carry_over_l_saved += saved_l
        pending_wastewater_l += fresh_water_l
        step_overflow_extra_l = 0.0

        groups={}
        for c in s.get("chemicals") or []:
            name=(c.get("chemical") or c.get("supplier") or c.get("company") or "").strip()
            if not name: continue
            dm=max(N(c.get("dose_min"),0),0)
            cb=N(c.get("begin_c"),cur)
            key=(round(fill_end+dm,3), round(cb,3), round(dm,3))
            groups.setdefault(key, []).append(c)

        for key in sorted(groups):
            target, cb, dm = key
            chems = groups[key]
            if abs(cb-cur) > .01:
                if cb > cur:
                    heating_kcal_raw += step_nominal_bath_l(s) * (cb - cur)
                rate = max(N(s.get("heating_slope") if cb>cur else s.get("cooling_gradient"),1), .1)
                segment_time = abs(cb-cur)/rate
                t += segment_time
                active_electric_time_min += segment_time
                cur=cb
                add_point(t, cur)
            if t < target:
                # Time between filling end and dose start is active machine time.
                # Filling itself is excluded, but dose_min waiting/circulation interval is included.
                active_electric_time_min += (target - t)
                t = target
                add_point(t, cur)
            labels=[]
            for c in chems:
                lab = letters[letter_i % len(letters)]
                letter_i += 1
                labels.append(lab)
                chemical_rows.append({
                    "label": lab, "supplier": c.get("supplier",""), "chemical": c.get("chemical",""),
                    "amount": N(c.get("amount"),0), "unit": c.get("unit","g/l"), "price": N(c.get("price"),0)
                })
                legend.append(f"{lab}: {c.get('chemical') or c.get('supplier') or 'Chemical'}")
            dose_time=max([N(c.get("dose_time"),0) for c in chems] or [0])
            end_temp=N(chems[-1].get("final_c"), cb) if chems else cb
            events.append({"type":"chemical_group","x":round(t,3),"y":round(cur,3),"labels":labels,"dose_time":round(dose_time,3),"y_end":round(end_temp,3)})
            if dose_time>0:
                if end_temp > cur:
                    heating_kcal_raw += step_nominal_bath_l(s) * (end_temp - cur)
                t += dose_time
                active_electric_time_min += dose_time
                cur = end_temp
                add_point(t, cur)
            circ=max([N(c.get("circulation_time"),0) for c in chems] or [0])
            if circ>0:
                t += circ
                active_electric_time_min += circ
                add_point(t, cur)

        final=N(s.get("final_temp"),cur)
        if abs(final-cur)>.01:
            if final > cur:
                heating_kcal_raw += step_nominal_bath_l(s) * (final - cur)
            rate=max(N(s.get("heating_slope") if final>cur else s.get("cooling_gradient"),1),.1)
            segment_time = abs(final-cur)/rate
            t += segment_time
            active_electric_time_min += segment_time
            cur=final
            add_point(t,cur)
        dwell=N(s.get("dwelling_time"),0)
        if dwell>0:
            t+=dwell
            active_electric_time_min += dwell
            add_point(t,cur)
        cool=N(s.get("cooling_temp"),cur)
        if cool<cur:
            rate=max(N(s.get("cooling_gradient"),1),.1)
            segment_time = abs(cur-cool)/rate
            t += segment_time
            active_electric_time_min += segment_time
            cur=cool
            add_point(t,cur)
        overflow=N(s.get("overflow_time"),0)
        if overflow>0:
            start_t = t
            base_temp = cur
            extra_l = overflow_extra_liters(s, overflow)
            overflow_extra_water_l += extra_l
            step_overflow_extra_l += extra_l
            pending_wastewater_l += extra_l
            active_electric_time_min += overflow
            events.append({"type":"overflow","x1":round(start_t,3),"x2":round(start_t+overflow,3),"y":round(base_temp,3),"extra_water_l":round(extra_l,2)})
            # Overflow is shown as a tighter ±10°C sine wave through the flow time.
            points = max(32, int(overflow * 6))
            for i in range(1, points + 1):
                tt = start_t + overflow * i / points
                yy = base_temp + 10.0 * math.sin(2 * math.pi * i / max(points/6, 1))
                add_point(tt, yy)
            t = start_t + overflow
            cur = base_temp
            add_point(t, cur)
        if s.get("drain"):
            # Drain creates a visual break in the temperature profile.
            # Do NOT connect previous step temperature to the next step beginning temperature.
            events.append({"type":"drain","x":round(t,3),"y":round(cur,3)})
            x.append(None); y.append(None)
            # Wastewater is counted only on drained steps and is limited to water actually taken.
            # If previous steps continued without drain, their taken water is discharged here.
            wastewater_l += max(pending_wastewater_l, 0.0)
            pending_wastewater_l = 0.0
            current_bath_l = 0.0
            t += drain_time
            cur = None
            previous_step_drained = True
        else:
            # No drain means no wastewater for this step; the bath remains for the next step.
            previous_step_drained = False

    total_water_l = base_process_water_l + overflow_extra_water_l
    water_m3 = total_water_l / 1000.0
    wastewater_m3 = wastewater_l / 1000.0
    total_time=t
    electricity_kwh = total_power * (active_electric_time_min/60)
    electricity_cost = electricity_kwh * N(utilities.get("electric_unit_price"), .15)
    water_cost = water_m3 * N(utilities.get("water_unit_price"), 1)
    waste_cost = wastewater_m3 * N(utilities.get("waste_water_unit_price"), .9)
    labour = (total_time/60) * N(utilities.get("hourly_wage"),0) * N(utilities.get("number_of_workers"),1) / max(N(utilities.get("number_of_machine"),1),1)
    heating_kcal = max(0, heating_kcal_raw) * (1+N(utilities.get("transfer_heat_loss"),0)/100)
    heating_consumption = heating_kcal / max(N(utilities.get("heating_capacity"),8250),1)
    heating_cost = heating_consumption * N(utilities.get("natural_gas_unit_price"),1.2)
    chemical_cost = 0.0

for r in chemical_rows:
    amt = r["amount"]

    # % uses fabric kg
    if r["unit"] == "%":
        chemical_cost += (fabric * amt / 100) * r["price"]

    # g/L uses step bath liters
    else:
        step_water_l = (
            r.get("bath_liters")
            or r.get("water_l")
            or r.get("step_water_l")
            or base_process_water_l
        )

        chemical_cost += (step_water_l * amt / 1000) * r["price"]
    total_cost = chemical_cost + electricity_cost + heating_cost + water_cost + waste_cost + labour
    electric_co2 = electricity_kwh * 0.42
    heating_co2 = heating_consumption * 2.02
    total_co2 = electric_co2 + heating_co2
    currency = project.get("cost_currency","EUR")
    dashboard = {
        f"Total Cost / batch ({currency})": round(total_cost,2),
        "Total Cost": round(total_cost,2),
        "Total Cost / kg": round(total_cost/fabric,3),
        "Total Time (min)": round(total_time,1),
        "Electricity (kWh/batch)": round(electricity_kwh,2),
        "Electricity Cost / batch": round(electricity_cost,2),
        "Heating Energy (kcal/batch)": round(heating_kcal,0),
        "Heating Consumption (Sm³/batch)": round(heating_consumption,2),
        "Heating Cost / batch": round(heating_cost,2),
        "Water Cost / batch": round(water_cost,2),
        "Waste Water Cost / batch": round(waste_cost,2),
        "Total CO2 / kg (g)": round(total_co2*1000/fabric,2),
        "Water L / kg": round((water_m3*1000)/fabric,2),
        "Total Water L / batch": round(water_m3*1000,2),
        "Wastewater L / batch": round(wastewater_l,2),
        "Wastewater L / kg": round(wastewater_l/fabric,2),
        "Active Electricity Time (min)": round(active_electric_time_min,1),
        "Base Process Water L / batch": round(base_process_water_l,2),
        "Carry Over Saved Water L / batch": round(carry_over_l_saved,2),
        "Overflow Extra Water L / batch": round(overflow_extra_water_l,2),
        "Energy kWh / kg": round(electricity_kwh/fabric,3),
    }

    return {
        "x": x,
        "y": y,
        "events": events,
        "chemical_rows": chemical_rows,
        "chemical_legend": legend,
        "dashboard": dashboard,
        "total_time": round(total_time, 1)
    }

@app.post("/api/calculate")
async def calculate(req: Request):
    return calc(await req.json())

def create_chart_png(project, out_path):
    data=calc(project)
    xs=data["x"]; ys=data["y"]
    fig, ax = plt.subplots(figsize=(14,5), dpi=180)
    curx=[]; cury=[]
    for a,b in zip(xs,ys):
        if a is None or b is None:
            if curx:
                ax.plot(curx, cury, linewidth=2.5)
                curx=[]; cury=[]
        else:
            curx.append(a); cury.append(b)
    if curx: ax.plot(curx, cury, linewidth=2.5)

    max_y=max([v for v in ys if isinstance(v,(int,float))] or [100])
    ax.set_ylim(0, max(110, math.ceil((max_y+10)/10)*10))
    ax.set_xlim(0, max([v for v in xs if isinstance(v,(int,float))] or [10]) * 1.03)
    ax.set_title("Dyeing Process Temperature Profile", fontsize=15, weight="bold")
    ax.set_xlabel("Time (min)", weight="bold")
    ax.set_ylabel("Temperature (°C)", weight="bold")
    ax.grid(True, linestyle="--", alpha=.35)
    for e in data["events"]:
        if e["type"]=="chemical_group":
            x=e["x"]; y=e["y"]; label=",".join(e.get("labels",[]))
            if e.get("dose_time",0)>1:
                ax.plot([x, x+e["dose_time"], x+e["dose_time"], x], [y, e.get("y_end",y), max(y,e.get("y_end",y))+8, y], color="black", linewidth=1.5)
                ax.text(x+e["dose_time"]/2, max(y,e.get("y_end",y))+10, label, ha="center", fontsize=8, weight="bold")
            else:
                ax.annotate(label, xy=(x,y), xytext=(x, y+18), ha="center", arrowprops=dict(arrowstyle="-|>", lw=1.2), fontsize=8, weight="bold")
        elif e["type"]=="drain":
            ax.annotate("Drain", xy=(e["x"], e["y"]), xytext=(e["x"], max(5, e["y"]-35)), ha="center", arrowprops=dict(arrowstyle="-|>", lw=1.8, linestyle="--"), fontsize=8, weight="bold")
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)

# ---------- PPT helpers ----------
def add_textbox(slide, text, x, y, w, h, font=14, bold=False, color=(20,36,60), align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=shape.text_frame
    tf.clear()
    tf.margin_left=Pt(6); tf.margin_right=Pt(6); tf.margin_top=Pt(3); tf.margin_bottom=Pt(3)
    p=tf.paragraphs[0]
    p.text=str(text)
    if align: p.alignment=align
    r=p.runs[0]
    r.font.size=Pt(font); r.font.bold=bold; r.font.color.rgb=RGBColor(*color)
    return shape

def add_rect(slide, x,y,w,h, fill=(255,255,255), line=(220,228,238), radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb=RGBColor(*fill)
    shape.line.color.rgb=RGBColor(*line)
    return shape

def set_bg(slide, color=(246,249,252)):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*color)

def header(slide, title, subtitle="DyeFlow RS Executive Report"):
    add_rect(slide,0,0,13.333,0.78,fill=(8,55,99),line=(8,55,99))
    add_textbox(slide,title,0.35,0.12,7.8,0.35,20,True,(255,255,255))
    add_textbox(slide,subtitle,0.36,0.48,7.8,0.2,8,False,(220,235,250))
    add_rect(slide,11.52,0.14,1.35,0.42,fill=(255,122,0),line=(255,122,0),radius=True)
    add_textbox(slide,"DyeFlow RS",11.56,0.23,1.28,0.18,9,True,(255,255,255),PP_ALIGN.CENTER)
    img=STATIC/"corner_image.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(12.75), Inches(0.08), width=Inches(0.42), height=Inches(0.55))

def card(slide, title, value, x,y,w,h):
    add_rect(slide,x,y,w,h,fill=(255,255,255),line=(214,226,238),radius=True)
    add_textbox(slide,title,x+0.08,y+0.08,w-0.16,0.2,8,True,(90,110,132))
    add_textbox(slide,value,x+0.08,y+0.33,w-0.16,h-0.4,16,True,(8,55,99))

def build_ppt(project, out_path):
    data=calc(project)
    d=data["dashboard"]
    chart_path=GENERATED/f"chart_{int(time.time()*1000)}.png"
    create_chart_png(project, chart_path)
    prs=Presentation()
    prs.slide_width=Inches(13.333)
    prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]

    # Slide 1
    s=prs.slides.add_slide(blank); set_bg(s); header(s,"Executive Inputs","Basic Inputs + Machine Parameters")
    add_textbox(s, project.get("project_name","Project"),0.45,0.95,6,0.3,17,True,(8,55,99))
    add_textbox(s, f"Company: {project.get('company_name','-')}   |   Process: {project.get('process_type','-')}",0.45,1.3,8,0.25,10,False,(80,95,115))
    basic=[("Fabric kg",project.get("fabric_kg",0)),("Flote",project.get("flote",0)),("Heating",project.get("utilities",{}).get("heating_source","-")),("Currency",project.get("cost_currency","EUR"))]
    mach=project.get("machine",{})
    machine=[("Machine",mach.get("machine_name","-")),("Capacity kg",mach.get("capacity_kg",0)),("Drain min",mach.get("drain_time_min",0)),("Total Power kW", round(N(mach.get("circulation_pump_power"))*N(mach.get("pump_ratio"),1)+N(mach.get("number_of_reel"))*N(mach.get("reel_power"))+N(mach.get("fan_power"))*N(mach.get("fan_ratio"),1),2))]
    add_textbox(s,"Basic Inputs",0.55,1.85,3,0.25,14,True,(8,55,99))
    add_textbox(s,"Machine",6.9,1.85,3,0.25,14,True,(8,55,99))
    for i,(k,v) in enumerate(basic): card(s,k,str(v),0.55+(i%2)*2.95,2.25+(i//2)*1.08,2.65,0.86)
    for i,(k,v) in enumerate(machine): card(s,k,str(v),6.9+(i%2)*2.95,2.25+(i//2)*1.08,2.65,0.86)
    card(s,"Total Cost / kg",str(d.get("Total Cost / kg",0)),0.55,5.0,2.65,0.9)
    card(s,"CO₂ / kg",f'{d.get("Total CO₂ / kg (g)",0)} g',3.5,5.0,2.65,0.9)
    card(s,"Total Time",f'{d.get("Total Time (min)",0)} min',6.45,5.0,2.65,0.9)
    card(s,"Water L/kg",str(d.get("Water L / kg",0)),9.4,5.0,2.65,0.9)

    # Slide 2
    s=prs.slides.add_slide(blank); set_bg(s); header(s,"Process Graph","Temperature profile, dosing and drain events")
    s.shapes.add_picture(str(chart_path), Inches(0.55), Inches(1.02), width=Inches(12.25), height=Inches(4.65))
    legend="\n".join(data["chemical_legend"][:16]) or "No chemicals"
    add_rect(s,0.55,5.86,12.25,1.0,fill=(255,255,255),line=(214,226,238),radius=True)
    add_textbox(s,"Chemical Legend",0.75,6.02,2.0,0.22,10,True,(8,55,99))
    add_textbox(s,legend,2.2,5.98,10.3,0.72,8,False,(42,55,70))

    # Slide 3
    s=prs.slides.add_slide(blank); set_bg(s); header(s,"Dashboard Report","Cost, energy and carbon overview")
    metrics=[
        ("Total Cost",d.get("Total Cost",0)),("Chemical Cost",d.get("Chemical Cost / batch",0)),
        ("Electricity kWh",d.get("Electricity (kWh/batch)",0)),("Heating Sm³",d.get("Heating Consumption (Sm³/batch)",0)),
        ("Total CO₂ kg",d.get("Total CO₂ (kg/batch)",0)),("CO₂ g/kg",d.get("Total CO₂ / kg (g)",0)),
        ("Water L/kg",d.get("Water L / kg",0)),("Energy kWh/kg",d.get("Energy kWh / kg",0))
    ]
    for i,(k,v) in enumerate(metrics):
        card(s,k,str(v),0.55+(i%4)*3.05,1.1+(i//4)*1.08,2.7,0.88)
    # Simple bars
    costs=[("Chemical",d.get("Chemical Cost / batch",0)),("Electricity",d.get("Electricity Cost / batch",0)),("Heating",d.get("Heating Cost / batch",0)),("Water",d.get("Water Cost / batch",0)),("Waste",d.get("Waste Water Cost / batch",0)),("Labour",d.get("Labour Cost / batch",0))]
    maxc=max(max([float(x[1]) for x in costs] or [1]),1)
    add_textbox(s,"Cost Breakdown",0.75,3.65,3,0.25,13,True,(8,55,99))
    for i,(k,v) in enumerate(costs):
        yy=4.05+i*0.38
        add_textbox(s,k,0.75,yy,1.3,0.2,8,True,(55,70,90))
        add_rect(s,2.0,yy,4.2,0.18,fill=(232,238,245),line=(232,238,245),radius=True)
        add_rect(s,2.0,yy,4.2*float(v)/maxc,0.18,fill=(255,122,0),line=(255,122,0),radius=True)
        add_textbox(s,str(round(float(v),2)),6.3,yy-0.02,0.8,0.2,8,False,(55,70,90))
    add_textbox(s,"Executive Note",7.55,3.65,2.5,0.25,13,True,(8,55,99))
    note=f"Batch total time is {d.get('Total Time (min)',0)} minutes. Total cost per kg is {d.get('Total Cost / kg',0)} and carbon intensity is {d.get('Total CO₂ / kg (g)',0)} g CO₂/kg."
    add_rect(s,7.55,4.0,4.85,1.95,fill=(255,255,255),line=(214,226,238),radius=True)
    add_textbox(s,note,7.75,4.2,4.45,1.3,13,False,(42,55,70))

    # Slide 4
    s=prs.slides.add_slide(blank); set_bg(s); header(s,"Cost Per Kg Analysis","Per kg cost distribution")
    perkg=[("Chemical",d.get("Chemical Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001)),("Electricity",d.get("Electricity Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001)),("Heating",d.get("Heating Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001)),("Water",d.get("Water Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001)),("Waste",d.get("Waste Water Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001)),("Labour",d.get("Labour Cost / batch",0)/max(N(project.get("fabric_kg"),1),.001))]
    maxv=max(max([v for _,v in perkg] or [1]),1)
    card(s,"Total Cost / kg",str(d.get("Total Cost / kg",0)),0.75,1.2,2.8,1.0)
    for i,(k,v) in enumerate(perkg):
        yy=2.65+i*0.52
        add_textbox(s,k,0.9,yy,1.4,0.22,10,True,(55,70,90))
        add_rect(s,2.55,yy,7.0,0.24,fill=(232,238,245),line=(232,238,245),radius=True)
        add_rect(s,2.55,yy,7.0*v/maxv,0.24,fill=(8,55,99),line=(8,55,99),radius=True)
        add_textbox(s,f"{v:.3f}/kg",9.75,yy-0.02,1.4,0.22,10,True,(8,55,99))

    # Slide 5
    s=prs.slides.add_slide(blank); set_bg(s); header(s,"Carbon Footprint Per Kg","Carbon and resource intensity dashboard")
    card(s,"CO₂ / kg",f'{d.get("Total CO₂ / kg (g)",0)} g',0.75,1.25,2.8,1.0)
    card(s,"Total CO₂ / batch",f'{d.get("Total CO₂ (kg/batch)",0)} kg',3.85,1.25,2.8,1.0)
    card(s,"Energy / kg",f'{d.get("Energy kWh / kg",0)} kWh',6.95,1.25,2.8,1.0)
    card(s,"Water / kg",f'{d.get("Water L / kg",0)} L',10.05,1.25,2.8,1.0)
    carbon=[("Electricity CO₂",d.get("Electricity CO₂ (kg/batch)",0)),("Heating CO₂",d.get("Heating CO₂ (kg/batch)",0))]
    maxv=max(max([v for _,v in carbon] or [1]),1)
    add_textbox(s,"CO₂ Source Breakdown",0.9,3.0,3.0,0.25,14,True,(8,55,99))
    for i,(k,v) in enumerate(carbon):
        yy=3.55+i*0.72
        add_textbox(s,k,0.9,yy,1.8,0.25,11,True,(55,70,90))
        add_rect(s,2.9,yy,7.2,0.32,fill=(232,238,245),line=(232,238,245),radius=True)
        add_rect(s,2.9,yy,7.2*v/maxv,0.32,fill=(255,122,0),line=(255,122,0),radius=True)
        add_textbox(s,f"{v:.2f} kg",10.35,yy,1.0,0.25,11,True,(8,55,99))
    add_rect(s,0.9,5.35,11.55,0.75,fill=(255,255,255),line=(214,226,238),radius=True)
    add_textbox(s,"Interpretation: This slide shows the kg-based footprint indicators for management comparison between recipes, machines and process alternatives.",1.1,5.55,11.1,0.25,12,False,(42,55,70))
    prs.save(out_path)
    try: chart_path.unlink()
    except Exception: pass





@app.post("/api/export/csv")
async def export_csv(req: Request):
    project=await req.json(); data=calc(project)
    sio=io.StringIO(); w=csv.writer(sio)
    w.writerow(["Metric","Value"])
    for k,v in data["dashboard"].items(): w.writerow([k,v])
    return StreamingResponse(iter([sio.getvalue()]), media_type="text/csv", headers={"Content-Disposition":"attachment; filename=DyeFlow_RS.csv"})

@app.post("/api/export/report-html")
async def export_html(req: Request):
    payload=await req.json()
    project = payload.get("project", payload) if isinstance(payload, dict) else payload
    chart_png = payload.get("chart_png") if isinstance(payload, dict) else None
    data=calc(project)
    d=data["dashboard"]
    machine=project.get("machine",{}) or {}
    utilities=project.get("utilities",{}) or {}
    def rows(items):
        return "".join(f"<tr><td>{k}</td><td>{v}</td></tr>" for k,v in items)
    basic_rows=rows([
        ("Project Name", project.get("project_name","")), ("Company", project.get("company_name","")),
        ("Process", project.get("process_type","")), ("Fabric kg", project.get("fabric_kg",0)),
        ("Flote", project.get("flote",0)), ("Heating Source", utilities.get("heating_source","")),
    ])
    machine_rows=rows([
        ("Machine", machine.get("machine_name","")), ("Capacity kg", machine.get("capacity_kg",0)),
        ("Drain Time min", machine.get("drain_time_min",0)), ("Pump kW", machine.get("circulation_pump_power",0)),
        ("Reel Power", machine.get("reel_power",0)), ("Fan Power", machine.get("fan_power",0)),
    ])
    dash_rows="".join(f"<tr><td>{k}</td><td>{v}</td></tr>" for k,v in d.items())
    if not chart_png:
        out = V38_GENERATED / f"report_chart_{int(time.time())}.png"
        v46_render_process_graph(project, out)
        chart_png = "data:image/png;base64," + base64.b64encode(out.read_bytes()).decode("ascii")
    html=f"""<!doctype html><html><head><meta charset='utf-8'><title>DyeFlow RS Report</title>
    <style>body{{font-family:Arial;padding:28px;color:#102a48}}h1,h2{{color:#083763}}table{{border-collapse:collapse;width:100%;margin:10px 0 20px}}td,th{{border:1px solid #d7e2ee;padding:8px}}th{{background:#083763;color:white}}.chart{{width:100%;border:1px solid #d7e2ee;margin:8px 0 22px}}</style></head>
    <body><h1>DyeFlow RS Report</h1><h2>Basic Inputs</h2><table><tr><th>Field</th><th>Value</th></tr>{basic_rows}</table>
    <h2>Machine Parameters</h2><table><tr><th>Field</th><th>Value</th></tr>{machine_rows}</table>
    <h2>Process Graph</h2><img class='chart' src='{chart_png}'/>
    <h2>Dashboard</h2><table><tr><th>Metric</th><th>Value</th></tr>{dash_rows}</table></body></html>"""
    return HTMLResponse(html)

@app.post("/api/export/package")
async def export_package(req: Request):
    project=await req.json()
    tmp=GENERATED/f"package_{int(time.time())}.zip"
    GENERATED.mkdir(exist_ok=True)
    ppt=GENERATED/"DyeFlow_RS.pptx"; chart=GENERATED/"DyeFlow_RS_Chart.png"
    build_ppt(project,ppt); create_chart_png(project,chart)
    with zipfile.ZipFile(tmp,"w",zipfile.ZIP_DEFLATED) as z:
        z.write(ppt,"DyeFlow_RS.pptx"); z.write(chart,"DyeFlow_RS_Chart.png")
        z.writestr("project.json",json.dumps(project,ensure_ascii=False,indent=2))
    return FileResponse(tmp, media_type="application/zip", filename="DyeFlow_RS_Package.zip")




@app.post("/api/compare")
async def compare(req: Request):
    data=await req.json(); ps=data.get("projects",[])
    if len(ps)<2: return {"rows":[]}
    c1=calc(ps[0])["dashboard"]; c2=calc(ps[1])["dashboard"]
    keys=["Total Cost / kg","Total CO₂ / kg (g)","Total Time (min)","Water L / kg","Energy kWh / kg"]
    rows=[]
    for k in keys:
        a=N(c1.get(k),0); b=N(c2.get(k),0)
        if abs(a-b) <= 0.005:
            diff=0.0; advantage="Equal"; status="equal"
        else:
            diff=((b-a)/a*100) if abs(a)>0.005 else (100.0 if b>a else -100.0)
            advantage="Project 1" if a<b else "Project 2"
            status="better_p1" if a<b else "better_p2"
        rows.append({"metric":k,"project1":a,"project2":b,"diff":round(diff,1),"advantage":advantage,"status":status})
    return {"project1":ps[0].get("project_name","Project 1"),"project2":ps[1].get("project_name","Project 2"),"rows":rows,"insights":[],"winner":""}

# ---------- Project files ----------
def _safe_name(name):
    return "".join(ch if ch.isalnum() or ch in "._- " else "_" for ch in name).strip() or "project"

@app.post("/api/save-project")
async def save_project(req: Request):
    u=_user_from_request(req)
    if not u: raise HTTPException(401,"Login required")
    if not u.get("can_save",True): raise HTTPException(403,"Save permission denied")
    data=await req.json()
    userdir=PROJECTS/u["username"]; userdir.mkdir(exist_ok=True)
    fname=_safe_name(data.get("project_name","DyeFlow_Project"))+".json"
    (userdir/fname).write_text(json.dumps(data,ensure_ascii=False,indent=2),encoding="utf-8")
    return {"file":fname}

@app.get("/api/projects")
def projects(req: Request):
    u=_user_from_request(req)
    if not u: raise HTTPException(401,"Login required")
    userdir=PROJECTS/u["username"]; userdir.mkdir(exist_ok=True)
    out=[]
    for p in sorted(userdir.glob("*.json"), key=lambda x:x.stat().st_mtime, reverse=True):
        try:
            d=json.loads(p.read_text(encoding="utf-8"))
        except Exception: d={}
        out.append({"file":p.name,"updated":datetime.fromtimestamp(p.stat().st_mtime).strftime("%Y-%m-%d %H:%M"),"project_name":d.get("project_name",""),"company_name":d.get("company_name","")})
    return {"projects":out}

@app.get("/api/load-project/{fname}")
def load_project(fname: str, req: Request):
    u=_user_from_request(req)
    if not u: raise HTTPException(401,"Login required")
    p=PROJECTS/u["username"]/_safe_name(fname)
    if not p.exists(): raise HTTPException(404,"Project not found")
    return json.loads(p.read_text(encoding="utf-8"))

@app.post("/api/upload-template")
async def upload_template(file: UploadFile = File(...)):
    dest=UPLOADS/_safe_name(file.filename)
    with dest.open("wb") as f:
        shutil.copyfileobj(file.file, f)
    return {"file":dest.name, "note":"Template upload saved. Default v28 export uses DyeFlow RS premium theme."}



# =========================
# v47 Stable Base: Premium Report + Compare Engine (v46 core preserved)
# =========================
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

V38_GENERATED = APP_DIR / "generated_ppt"
V38_GENERATED.mkdir(exist_ok=True)

def v46_num(x, default=0.0):
    try:
        if x is None or x == "":
            return float(default)
        return float(x)
    except Exception:
        return float(default)

def v46_calc(project):
    # Use existing project calculation engine if available; preserve all old keys.
    try:
        base_calc = calc(project)
    except Exception:
        base_calc = {"dashboard": {}, "x": [], "y": [], "events": [], "chemical_rows": [], "chemical_legend": []}

    dash = dict(base_calc.get("dashboard") or {})
    fabric = max(v46_num(project.get("fabric_kg"), 1), 0.001)

    # Backward-compatible aliases.
    aliases = {
        "Heating Cost / batch": ["Heating Cost / batch", "Heating Cost", "Heating"],
        "Electricity Cost / batch": ["Electricity Cost / batch", "Electricity Cost", "Electricity"],
        "Chemical Cost / batch": ["Chemical Cost / batch", "Chemical Cost", "Chemical"],
        "Water Cost / batch": ["Water Cost / batch", "Water Cost", "Water"],
        "Waste Water Cost / batch": ["Waste Water Cost / batch", "Waste Water Cost", "Wastewater Cost"],
        "Labour Cost / batch": ["Labour Cost / batch", "Labour Cost", "Labor Cost"],
        "Total Cost / batch": ["Total Cost / batch", "Total Cost"],
        "Electricity CO₂ (kg/batch)": ["Electricity CO₂ (kg/batch)", "Electricity CO2 (kg/batch)", "Electricity CO₂"],
        "Heating CO₂ (kg/batch)": ["Heating CO₂ (kg/batch)", "Heating CO2 (kg/batch)", "Heating CO₂"],
        "Total CO₂ (kg/batch)": ["Total CO₂ (kg/batch)", "Total CO2 (kg/batch)", "Total CO₂"],
        "Total CO₂ / kg (g)": ["Total CO₂ / kg (g)", "Total CO2 / kg (g)", "CO₂ / kg"],
        "Water L / kg": ["Water L / kg"],
        "Energy kWh / kg": ["Energy kWh / kg"],
        "Total Time (min)": ["Total Time (min)"],
    }

    def pick(*keys):
        for k in keys:
            if k in dash:
                return v46_num(dash.get(k), 0)
        return 0.0

    normalized = {}
    for target, keys in aliases.items():
        normalized[target] = pick(*keys)

    # If total cost not available, sum components.
    if not normalized["Total Cost / batch"]:
        normalized["Total Cost / batch"] = sum(normalized[k] for k in [
            "Heating Cost / batch", "Electricity Cost / batch", "Chemical Cost / batch",
            "Water Cost / batch", "Waste Water Cost / batch", "Labour Cost / batch"
        ])

    for k in ["Heating Cost / batch", "Electricity Cost / batch", "Chemical Cost / batch",
              "Water Cost / batch", "Waste Water Cost / batch", "Labour Cost / batch",
              "Total Cost / batch"]:
        normalized[k.replace(" / batch", " / kg")] = normalized[k] / fabric

    if not normalized["Total CO₂ / kg (g)"] and normalized["Total CO₂ (kg/batch)"]:
        normalized["Total CO₂ / kg (g)"] = normalized["Total CO₂ (kg/batch)"] * 1000 / fabric

    normalized["Electricity CO₂ / kg (g)"] = normalized["Electricity CO₂ (kg/batch)"] * 1000 / fabric if fabric else 0
    normalized["Heating CO₂ / kg (g)"] = normalized["Heating CO₂ (kg/batch)"] * 1000 / fabric if fabric else 0

    dash.update({k: round(v, 3) for k, v in normalized.items()})
    base_calc["dashboard"] = dash
    return base_calc

def v46_points_from_project(project):
    # Recreate point/events with strict drain break.
    steps = project.get("steps") or []
    machine = project.get("machine") or {}
    drain_time = v46_num(machine.get("drain_time_min"), 5)
    x, y, events, chemical_rows, legend = [], [], [], [], []
    t = 0.0
    letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    li = 0
    first = True
    gap = False

    for si, s in enumerate(steps, start=1):
        begin = v46_num(s.get("beginning_temp"), 25)
        cur = begin

        if first:
            x.append(t); y.append(cur); first = False
        else:
            if gap:
                x.append(None); y.append(None)
                gap = False
            x.append(t); y.append(cur)

        fill = v46_num(s.get("filling_time"), 0)
        if fill > 0:
            t += fill
            x.append(t); y.append(cur)

        step_start = t - fill
        fill_end = step_start + fill
        groups = {}
        for c in s.get("chemicals", []) or []:
            if not any(str(c.get(k, "")).strip() for k in ["chemical", "supplier", "company"]):
                continue
            cb = v46_num(c.get("begin_c"), cur)
            dm = max(v46_num(c.get("dose_min"), 0), 0)
            target = fill_end + dm
            key = (round(target, 3), round(cb, 3))
            groups.setdefault(key, []).append(c)

        for (target, cb), chems in sorted(groups.items()):
            if abs(cb - cur) > 0.01:
                rate = max(v46_num(s.get("heating_slope") if cb > cur else s.get("cooling_gradient"), 1), 0.1)
                t += abs(cb-cur) / rate
                cur = cb
                x.append(t); y.append(cur)
            if t < target:
                t = target
                x.append(t); y.append(cur)

            labels = []
            dose_time = max([v46_num(c.get("dose_time"), 0) for c in chems] or [0])
            end_temp = cb
            for c in chems:
                lab = letters[li] if li < len(letters) else f"X{li+1}"
                li += 1
                labels.append(lab)
                end_temp = v46_num(c.get("final_c"), cb)
                chem_name = c.get("chemical") or c.get("supplier") or "Chemical"
                legend.append(f"{lab}: {chem_name}, {v46_num(c.get('amount'),0):g} {c.get('unit','g/l')}")
                chemical_rows.append({
                    "step": si, "letter": lab, "supplier": c.get("supplier",""),
                    "chemical": c.get("chemical",""), "amount": v46_num(c.get("amount"),0),
                    "unit": c.get("unit","g/l"), "dose_time": dose_time,
                    "time": round(t,2), "begin_temp": cb, "final_temp": end_temp
                })
            events.append({"type": "chemical_group", "x": round(t,3), "y": round(cur,3), "labels": labels, "dose_time": round(dose_time,3), "y_end": round(end_temp,3)})
            if dose_time > 0:
                t += dose_time
                active_electric_time_min += dose_time
                cur = end_temp
                x.append(t); y.append(cur)

            circ = max([v46_num(c.get("circulation_time"), 0) for c in chems] or [0])
            if circ > 0:
                t += circ
                x.append(t); y.append(cur)

        final = v46_num(s.get("final_temp"), cur)
        if abs(final - cur) > 0.01:
            rate = max(v46_num(s.get("heating_slope") if final > cur else s.get("cooling_gradient"), 1), 0.1)
            t += abs(final-cur) / rate
            cur = final
            x.append(t); y.append(cur)

        dwell = v46_num(s.get("dwelling_time"), 0)
        if dwell > 0:
            t += dwell
            x.append(t); y.append(cur)

        cool = v46_num(s.get("cooling_temp"), cur)
        if cool < cur:
            rate = max(v46_num(s.get("cooling_gradient"), 1), 0.1)
            t += abs(cur-cool) / rate
            cur = cool
            x.append(t); y.append(cur)

        overflow = v46_num(s.get("overflow_time"), 0)
        if overflow > 0:
            start_t = t
            base_temp = cur
            events.append({"type":"overflow", "x1": round(start_t,3), "x2": round(start_t+overflow,3), "y": round(base_temp,3)})
            points = max(32, int(overflow * 6))
            for i in range(1, points + 1):
                tt = start_t + overflow * i / points
                yy = base_temp + 10.0 * math.sin(2 * math.pi * i / max(points/6, 1))
                x.append(tt); y.append(yy)
            t = start_t + overflow
            cur = base_temp
            x.append(t); y.append(cur)

        if s.get("drain"):
            events.append({"type": "drain", "x": round(t,3), "y": round(cur,3)})
            t += drain_time
            gap = True

    return {"x": x, "y": y, "events": events, "chemical_rows": chemical_rows, "chemical_legend": legend}

def v46_render_process_graph(project, out_path, title="Dyeing Process Temperature Profile", width=13.5, height=5.0):
    data = v46_points_from_project(project)
    xs, ys = data["x"], data["y"]
    fig, ax = plt.subplots(figsize=(width, height), dpi=190)

    # Red temperature profile with strict breaks.
    segx, segy = [], []
    for a, b in zip(xs, ys):
        if a is None or b is None:
            if segx:
                ax.plot(segx, segy, color="#c00000", linewidth=2.8, solid_capstyle="round")
                segx, segy = [], []
        else:
            segx.append(a); segy.append(b)
    if segx:
        ax.plot(segx, segy, color="#c00000", linewidth=2.8, solid_capstyle="round")

    numeric_x = [v for v in xs if isinstance(v, (int, float))]
    numeric_y = [v for v in ys if isinstance(v, (int, float))]
    max_x = max(numeric_x or [10])
    max_y = max(numeric_y or [100])
    ax.set_xlim(0, max_x * 1.04 if max_x else 10)
    ax.set_ylim(0, max(110, int((max_y + 20) / 10) * 10))
    ax.set_title(title, fontsize=15, weight="bold", pad=12)
    ax.set_xlabel("Time (min)", fontsize=10, weight="bold")
    ax.set_ylabel("Temperature (°C)", fontsize=10, weight="bold")
    ax.grid(True, linestyle="--", linewidth=0.7, alpha=0.35)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    y_top = ax.get_ylim()[1]
    for e in data["events"]:
        if e.get("type") == "overflow":
            ax.axvspan(e.get("x1", 0), e.get("x2", 0), color="#dff3ff", alpha=0.20, hatch="///", edgecolor="#b8dceb", linewidth=0.0)
            ax.text((e.get("x1",0)+e.get("x2",0))/2, 5, "Overflow", ha="center", va="bottom", fontsize=7, weight="bold", color="#5d7f8e")
    for e in data["events"]:
        if e["type"] == "drain":
            # Downward drain arrow starting from temperature line.
            ax.annotate("Drain", xy=(e["x"], e["y"]), xytext=(e["x"], max(5, e["y"] - 35)),
                        ha="center", va="top", fontsize=8, weight="bold",
                        arrowprops=dict(arrowstyle="-|>", color="black", lw=1.8, linestyle="--"))
        elif e["type"] == "chemical_group":
            label = ",".join(e.get("labels", []))
            dose_time = float(e.get("dose_time", 0) or 0)
            x0, y0 = e["x"], e["y"]
            y_end = e.get("y_end", y0)
            if dose_time > 1:
                # Solid black dosing triangle.
                tri_top = min(y_top - 6, max(y0, y_end) + 15)
                ax.plot([x0, x0 + dose_time, x0 + dose_time, x0],
                        [y0, y_end, tri_top, y0], color="black", linewidth=1.4)
                ax.text(x0 + dose_time / 2, tri_top + 3, label, ha="center", va="bottom",
                        fontsize=8, weight="bold")
            else:
                ax.annotate(label, xy=(x0, y0), xytext=(x0, min(y_top - 4, y0 + 24)),
                            ha="center", fontsize=8, weight="bold",
                            arrowprops=dict(arrowstyle="-|>", color="black", lw=1.2))

    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return data

def v46_add_text(slide, text, x, y, w, h, font=12, bold=False, color=(20,36,60), align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Pt(5); tf.margin_right = Pt(5); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    # Apply identical font formatting to every paragraph/run. This fixes the
    # chemical legend issue where only the first line kept the intended point size.
    for para in tf.paragraphs:
        if align:
            para.alignment = align
        for run in para.runs:
            run.font.size = Pt(font)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
    return shape

def v46_rect(slide, x, y, w, h, fill=(255,255,255), line=(214,226,238), radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    return shape

def v46_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(246,249,252)

def v46_header(slide, title, subtitle="DyeFlow RS Premium Report"):
    v46_rect(slide, 0, 0, 13.333, 0.72, fill=(8,55,99), line=(8,55,99), radius=False)
    v46_add_text(slide, title, 0.35, 0.11, 7.2, 0.28, 19, True, (255,255,255))
    v46_add_text(slide, subtitle, 0.36, 0.43, 8.2, 0.18, 8, False, (220,235,250))
    v46_rect(slide, 11.45, 0.14, 1.38, 0.40, fill=(255,122,0), line=(255,122,0), radius=True)
    v46_add_text(slide, "DyeFlow RS", 11.5, 0.24, 1.28, 0.14, 8, True, (255,255,255), PP_ALIGN.CENTER)
    img = APP_DIR / "static" / "corner_image.png"
    if img.exists():
        try:
            slide.shapes.add_picture(str(img), Inches(12.86), Inches(0.08), width=Inches(0.36), height=Inches(0.50))
        except Exception:
            pass

def v46_card(slide, title, value, x, y, w, h, accent=False):
    v46_rect(slide, x, y, w, h, fill=(255,255,255), line=(214,226,238), radius=True)
    v46_add_text(slide, title, x+0.07, y+0.07, w-0.14, 0.15, 7.6, True, (90,110,132))
    v46_add_text(slide, value, x+0.07, y+0.30, w-0.14, h-0.34, 14 if len(str(value)) < 10 else 11, True, (255,122,0) if accent else (8,55,99))



# ---------- v55 Presentation Polish Theme (based on v54) ----------
def v54_color(name):
    # v55: stronger contrast, more vivid enterprise navy/orange, darker body text.
    palette = {
        "navy": (0, 25, 64), "blue": (0, 68, 135), "accent": (235, 64, 0),
        "orange": (255, 118, 0), "soft": (242, 247, 252), "line": (202, 214, 228),
        "muted": (72, 86, 105), "text": (6, 22, 45), "green": (0, 150, 90),
        "red": (220, 30, 30), "pale_blue": (232, 243, 252), "ink": (3, 18, 42)
    }
    return palette.get(name, (0,0,0))

def v54_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245,248,252)

def v54_text(slide, text, x, y, w, h, font=11, bold=False, color=None, align=None):
    # v57_READABILITY_FINAL: PowerPoint export readability rule.
    # No generated PPT text should use a font smaller than 12 pt.
    font = max(float(font or 12), 12.0)
    color = color or v54_color("text")
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.text = str(text)
    if align: p.alignment = align
    for para in tf.paragraphs:
        if align: para.alignment = align
        for run in para.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)
    return shape

def v54_rect(slide, x, y, w, h, fill=(255,255,255), line=None, radius=True):
    line = line if line is not None else v54_color("line")
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    shape.line.width = Pt(0.7)
    return shape

def v54_header(slide, num, title, subtitle="", project=None):
    v54_text(slide, f"{num:02d}", 0.28, 0.22, 0.48, 0.42, 18, True, (255,255,255), PP_ALIGN.CENTER)
    v54_rect(slide, 0.20, 0.16, 0.62, 0.58, fill=v54_color("accent"), line=v54_color("accent"), radius=True)
    v54_text(slide, title, 1.02, 0.18, 6.2, 0.34, 21, True, v54_color("navy"))
    v54_text(slide, subtitle, 1.03, 0.56, 6.6, 0.18, 9.2, False, v54_color("muted"))
    v54_text(slide, "DyeFlow", 11.66, 0.20, 0.78, 0.22, 12, True, v54_color("navy"), PP_ALIGN.RIGHT)
    v54_rect(slide, 12.47, 0.19, 0.28, 0.23, fill=(255,255,255), line=v54_color("accent"), radius=True)
    v54_text(slide, "RS", 12.49, 0.22, 0.22, 0.11, 7, True, v54_color("accent"), PP_ALIGN.CENTER)
    if project:
        meta = f"{project.get('project_name','Project')}  |  {datetime.now().strftime('%d %b %Y')}"
        v54_text(slide, meta, 9.15, 0.58, 3.55, 0.16, 8.0, False, v54_color("muted"), PP_ALIGN.RIGHT)
    # subtle top rule
    v54_rect(slide, 0.2, 0.88, 12.9, 0.01, fill=(226,232,240), line=(226,232,240), radius=False)

def v54_footer(slide, page):
    v54_text(slide, "DyeFlow", 0.28, 7.02, 0.58, 0.18, 10, True, v54_color("navy"))
    v54_rect(slide, 0.88, 7.04, 0.22, 0.18, fill=(255,255,255), line=v54_color("accent"), radius=True)
    v54_text(slide, "RS", 0.895, 7.065, 0.17, 0.07, 5.8, True, v54_color("accent"), PP_ALIGN.CENTER)
    v54_text(slide, "Smart Process. Lower Cost. Better Future.", 5.25, 7.05, 3.0, 0.16, 6.8, False, v54_color("muted"), PP_ALIGN.CENTER)
    v54_text(slide, f"Page {page}", 12.22, 7.05, 0.66, 0.16, 6.8, False, v54_color("muted"), PP_ALIGN.RIGHT)

def v54_section_label(slide, title, x, y, w, dark=False):
    if dark:
        v54_rect(slide, x, y, w, 0.31, fill=v54_color("navy"), line=v54_color("navy"), radius=True)
        v54_text(slide, title.upper(), x+0.08, y+0.09, w-0.16, 0.11, 6.5, True, (255,255,255), PP_ALIGN.CENTER)
    else:
        v54_text(slide, "—  " + title, x, y, w, 0.16, 8, True, v54_color("accent"))

def v54_kpi(slide, title, value, unit, x, y, w, h, icon="●", accent=False):
    # v55: stronger KPI card contrast + larger numeric hierarchy
    v54_rect(slide, x, y, w, h, fill=(255,255,255), line=(190,205,222), radius=True)
    v54_rect(slide, x, y, 0.055, h, fill=v54_color("accent") if accent else v54_color("blue"), line=v54_color("accent") if accent else v54_color("blue"), radius=False)
    v54_text(slide, icon, x+0.10, y+0.12, 0.34, 0.25, 17, True, v54_color("accent") if accent else v54_color("blue"), PP_ALIGN.CENTER)
    v54_text(slide, title.upper(), x+0.42, y+0.13, w-0.52, 0.12, 7.2, True, v54_color("navy"), PP_ALIGN.CENTER)
    v54_text(slide, value, x+0.10, y+0.48, w-0.20, 0.25, 19, True, v54_color("ink"), PP_ALIGN.CENTER)
    v54_text(slide, unit, x+0.10, y+0.82, w-0.20, 0.12, 8.0, True, v54_color("text"), PP_ALIGN.CENTER)

def v54_info_panel(slide, title, rows, x, y, w, h):
    v54_rect(slide, x, y, w, h, fill=(255,255,255), line=(190,205,222), radius=True)
    v54_section_label(slide, title, x+0.12, y+0.12, w-0.24, dark=True)
    row_h = (h-0.62) / max(len(rows),1)
    yy = y+0.58
    for k,v in rows[:8]:
        v54_text(slide, str(k), x+0.18, yy+0.08, w*0.46, 0.12, 7.4, False, v54_color("muted"))
        v54_text(slide, str(v), x+w*0.50, yy+0.08, w*0.44, 0.12, 7.5, True, v54_color("text"))
        if yy < y+h-0.30:
            v54_rect(slide, x+0.14, yy+row_h-0.02, w-0.28, 0.004, fill=(225,232,241), line=(225,232,241), radius=False)
        yy += row_h

def v54_legend_panel(slide, legend, x, y, w, h, title="CHEMICAL LEGEND"):
    """v56 legend fix: no overlapping chemical items.

    The compare PPT uses a very wide but shallow legend panel. The old adaptive
    vertical spacing compressed 5-6 chemicals into less than 0.5 inch, causing
    stacked color chips and overlapping names. This version switches shallow
    wide panels to a wrapped grid layout and keeps tall panels as a vertical
    list. Calculation/graph engines are untouched.
    """
    v54_rect(slide, x, y, w, h, fill=(255,255,255), line=v54_color("line"), radius=True)
    header_h = 0.30 if h >= 1.0 else 0.24
    v54_rect(slide, x+0.06, y+0.08, w-0.12, header_h, fill=v54_color("navy"), line=v54_color("navy"), radius=True)
    v54_text(slide, title, x+0.10, y+0.18 if h>=1.0 else y+0.15, w-0.20, 0.09, 6.8 if h < 1.0 else 7.6, True, (255,255,255), PP_ALIGN.CENTER)
    colors=[(0,153,80),(255,144,0),(0,103,184),(126,87,194),(0,45,105),(220,70,140),(20,184,166),(120,85,50)]
    items=list(legend or [])
    if not items:
        v54_text(slide, "No chemical defined", x+0.15, y+header_h+0.30, w-0.30, 0.12, 7.2, False, v54_color("muted"))
        return

    def split_item(item, idx):
        if ":" in str(item):
            label, name = (str(item).split(":",1)+[""])[:2]
        else:
            label, name = chr(65+idx), str(item)
        return label.strip()[:3], (name.strip() or str(item)).strip()

    # Shallow compare legend panels: wrap items into 2 rows x 3 columns instead
    # of forcing a vertical list. This is the key fix for overlapping chemicals.
    if h < 1.25 and w >= 3.2:
        max_items = min(len(items), 6)
        cols = 3 if max_items > 3 else max_items
        rows = int(math.ceil(max_items / max(cols, 1)))
        content_top = y + header_h + 0.20
        content_h = max(0.34, h - header_h - 0.30)
        row_h = content_h / max(rows, 1)
        col_w = (w - 0.34) / max(cols, 1)
        box = min(0.18, max(0.13, row_h * 0.48))
        font = 6.9 if rows <= 2 else 6.2
        for i, item in enumerate(items[:max_items]):
            label, nm = split_item(item, i)
            r, cidx = divmod(i, cols)
            xx = x + 0.16 + cidx * col_w
            yy = content_top + r * row_h + max(0.00, (row_h - box) / 2)
            if len(nm) > 19:
                nm = nm[:18] + "…"
            col = colors[i % len(colors)]
            v54_rect(slide, xx, yy, box, box, fill=col, line=col, radius=True)
            v54_text(slide, label, xx+0.012, yy+box*0.29, box-0.024, 0.06, 4.8, True, (255,255,255), PP_ALIGN.CENTER)
            v54_text(slide, nm, xx+box+0.07, yy+box*0.20, col_w-box-0.10, 0.09, font, True, v54_color("text"))
        if len(items) > max_items:
            v54_text(slide, f"+{len(items)-max_items}", x+w-0.36, y+h-0.20, 0.22, 0.09, 5.5, True, v54_color("accent"), PP_ALIGN.RIGHT)
        return

    # Normal/tall legend panels: stable vertical list with a minimum row height.
    max_items = min(len(items), 12 if h >= 2.5 else 8)
    content_top = y + header_h + 0.24
    content_bottom = y + h - 0.22
    available = max(0.40, content_bottom - content_top)
    step = max(0.20, min(0.36, available / max(max_items, 1)))
    font = 7.4 if h >= 1.4 else 6.7
    yy = content_top
    for i,item in enumerate(items[:max_items]):
        label, nm = split_item(item, i)
        c=colors[i%len(colors)]
        box = min(0.22, max(0.16, step*0.62))
        if len(nm)>30 and h<1.4: nm=nm[:29]+"…"
        elif len(nm)>42: nm=nm[:41]+"…"
        v54_rect(slide, x+0.14, yy, box, box, fill=c, line=c, radius=True)
        v54_text(slide, label, x+0.155, yy+box*0.27, box-0.03, 0.06, 5.4, True, (255,255,255), PP_ALIGN.CENTER)
        v54_text(slide, nm, x+0.14+box+0.10, yy+box*0.20, w-0.30-box, 0.10, font, True if h<1.2 else False, v54_color("text"))
        yy += step
    if len(items) > max_items:
        v54_text(slide, f"+{len(items)-max_items} more chemicals", x+0.16, y+h-0.20, w-0.3, 0.10, 6.2, True, v54_color("accent"))

def v54_progress(slide, label, value, max_value, x, y, w, color=None, amount_text=""):
    color = color or v54_color("blue")
    v54_text(slide, label, x, y, 1.25, 0.12, 7.5, True, v54_color("text"))
    v54_rect(slide, x+1.35, y+0.02, w-2.25, 0.08, fill=(236,241,247), line=(236,241,247), radius=True)
    bw = (w-2.25) * (float(value)/max(float(max_value), 0.0001))
    v54_rect(slide, x+1.35, y+0.02, max(0.02,bw), 0.08, fill=color, line=color, radius=True)
    v54_text(slide, amount_text or f"{value:.2f}", x+w-0.76, y-0.01, 0.74, 0.12, 7.0, True, v54_color("muted"), PP_ALIGN.RIGHT)

def v54_fmt(v, nd=2):
    try:
        v=float(v)
        if abs(v) >= 100: return f"{v:.0f}"
        if abs(v) >= 10: return f"{v:.1f}"
        return f"{v:.{nd}f}"
    except Exception:
        return str(v)

def v54_cost_components(d):
    return [
        ("Heating", d.get("Heating Cost / batch",0), (230,57,70)),
        ("Electricity", d.get("Electricity Cost / batch",0), (0,92,184)),
        ("Chemical", d.get("Chemical Cost / batch",0), (245,158,11)),
        ("Water", d.get("Water Cost / batch",0), (0,150,136)),
        ("Waste", d.get("Waste Water Cost / batch",0), (126,87,194)),
        ("Labour", d.get("Labour Cost / batch",0), (5,28,58)),
    ]

def v54_donut(slide, components, x, y, size=1.0):
    total=sum(float(v or 0) for _,v,_ in components) or 1
    start=0
    # python-pptx pie segments are cumbersome; use a clean stacked mini legend + center disk
    v54_rect(slide, x, y, size, size, fill=(255,255,255), line=v54_color("line"), radius=True)
    v54_text(slide, "Cost\nStructure", x+0.08, y+0.33, size-0.16, 0.20, 8, True, v54_color("navy"), PP_ALIGN.CENTER)
    yy=y+size+0.10
    for lab,val,col in components[:6]:
        pct=100*float(val or 0)/total
        v54_rect(slide, x, yy, 0.10, 0.10, fill=col, line=col, radius=True)
        v54_text(slide, f"{lab}  {pct:.1f}%", x+0.15, yy-0.01, 1.05, 0.10, 5.6, False, v54_color("text"))
        yy += 0.17

def v54_build_single_ppt(project, out, chart_svg=None, chart_png=None):
    data = v46_calc(project)
    d = data.get("dashboard", {})
    graph_path = v46_png_dataurl_to_file(chart_png, "v54_single_graph") if chart_png else None
    if not graph_path:
        graph_path = V38_GENERATED / f"v54_single_graph_{int(time.time()*1000)}.png"
        try:
            v46_render_process_graph(project, graph_path)
        except Exception:
            create_chart_png(project, graph_path)
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    fabric=max(v46_num(project.get("fabric_kg"),1),0.001)
    m=project.get("machine",{}) or {}; u=project.get("utilities",{}) or {}
    total_power = round(v46_num(m.get("circulation_pump_power"))*v46_num(m.get("pump_ratio"),1)+v46_num(m.get("number_of_reel"))*v46_num(m.get("reel_power"))+v46_num(m.get("fan_power"))*v46_num(m.get("fan_ratio"),1),2)
    # Slide 1
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,1,"Executive Summary","DyeFlow RS",project)
    basic=[("Process",project.get("process_type","-")),("Fabric",f"{project.get('fabric_kg',0)} kg"),("Flote",project.get("flote",0)),("Heating",u.get("heating_source","-")),("Currency",project.get("cost_currency","EUR")),("Carry Over",f"{project.get('carry_over',0)} %"),("Fabric Status",project.get("fabric_status","-"))]
    mach=[("Machine",m.get("machine_name","-")),("Capacity",f"{m.get('capacity_kg',0)} kg"),("Drain Time",f"{m.get('drain_time_min',0)} min"),("Pump Power",f"{m.get('circulation_pump_power',0)} kW"),("Reel Power",f"{m.get('reel_power',0)} kW"),("Total Power",f"{total_power} kW")]
    v54_info_panel(s,"Project Information",basic,0.35,1.10,1.95,4.95)
    v54_info_panel(s,"Machine Information",mach,2.48,1.10,1.95,4.95)
    kx, ky, kw, kh = 4.60,1.02,1.65,1.45
    v54_kpi(s,"Total Cost",f"€{v54_fmt(d.get('Total Cost / batch',0),1)}",f"{v54_fmt(d.get('Total Cost / kg',0),3)} €/kg",kx,ky,kw,kh,"●",True)
    v54_kpi(s,"Total Time",v54_fmt(d.get("Total Time (min)",0),1),"min",kx+1.88,ky,kw,kh,"◷")
    v54_kpi(s,"Energy",v54_fmt(d.get("Energy kWh / kg",0),2),"kWh/kg",kx+3.76,ky,kw,kh,"⚡")
    v54_kpi(s,"Water",v54_fmt(d.get("Water L / kg",0),1),"L/kg",kx,ky+1.70,kw,kh,"◊")
    v54_kpi(s,"CO₂ Emission",v54_fmt(d.get("Total CO₂ (kg/batch)",0),1),f"kg/batch | {v54_fmt(d.get('Total CO₂ / kg (g)',0),1)} g/kg",kx+1.88,ky+1.70,kw,kh,"☁")
    v54_kpi(s,"Chemical Cost",f"€{v54_fmt(d.get('Chemical Cost / batch',0),1)}",f"{v54_fmt(d.get('Chemical Cost / kg',0),3)} €/kg",kx+3.76,ky+1.70,kw,kh,"⚗")
    v54_rect(s,0.35,6.15,12.45,0.48,fill=(255,255,255),line=v54_color("line"),radius=True)
    insight=f"This process has a total cost of €{v54_fmt(d.get('Total Cost / kg',0),3)} per kg with {v54_fmt(d.get('Total CO₂ / kg (g)',0),1)} g CO₂ emission and {v54_fmt(d.get('Water L / kg',0),1)} L water usage per kg."
    v54_text(s,"⚙  INSIGHT",0.55,6.31,1.0,0.10,7.4,True,v54_color("accent")); v54_text(s,insight,1.55,6.29,10.7,0.13,7.3,False,v54_color("text"))
    v54_footer(s,1)
    # Slide 2
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,2,"Process Timeline","Temperature & Chemical Dosing Profile",project)
    v54_rect(s,0.30,1.06,12.75,5.25,fill=(255,255,255),line=v54_color("line"),radius=True)
    # process ribbon
    stages=[("1. FILLING","◊"),("2. HEATING","♨"),("3. PROCESS","⟳"),("4. COOLING","↘"),("5. DRAIN","↓"),("6. EXTRACTION","✤")]
    sx=0.92
    for i,(lab,ico) in enumerate(stages):
        v54_text(s,ico,sx+i*1.85,1.18,0.34,0.22,13,False,v54_color("navy"),PP_ALIGN.CENTER)
        v54_text(s,lab,sx+0.37+i*1.85,1.22,1.05,0.12,6.9,True,v54_color("navy"),PP_ALIGN.CENTER)
        if i<5: v54_rect(s,sx+1.62+i*1.85,1.14,0.005,0.40,fill=(220,228,238),line=(220,228,238),radius=False)
    s.shapes.add_picture(str(graph_path), Inches(0.55), Inches(1.72), width=Inches(10.30), height=Inches(4.05))
    v54_legend_panel(s,data.get("chemical_legend",[]),11.05,1.90,1.72,3.25)
    v54_rect(s,0.35,6.43,12.50,0.46,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,f"Total Cost €{v54_fmt(d.get('Total Cost / batch',0),1)}",0.62,6.58,2.1,0.10,7.8,True,v54_color("navy"))
    v54_text(s,f"Time {v54_fmt(d.get('Total Time (min)',0),1)} min",3.0,6.58,2.0,0.10,7.8,True,v54_color("navy"))
    v54_text(s,f"Water {v54_fmt(d.get('Water L / kg',0),1)} L/kg",5.25,6.58,2.0,0.10,7.8,True,v54_color("navy"))
    v54_text(s,f"CO₂ {v54_fmt(d.get('Total CO₂ / kg (g)',0),1)} g/kg",7.55,6.58,2.0,0.10,7.8,True,v54_color("navy"))
    v54_text(s,f"Energy {v54_fmt(d.get('Energy kWh / kg',0),2)} kWh/kg",9.95,6.58,2.3,0.10,7.8,True,v54_color("navy"))
    v54_footer(s,2)
    # Slide 3 chemical operations
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,3,"Chemical Operations","Step by Step Chemical Dosing Summary",project)
    rows=data.get("chemical_rows",[])[:8]
    v54_rect(s,0.35,1.08,12.45,4.85,fill=(255,255,255),line=v54_color("line"),radius=True)
    headers=[("STEP",0.62), ("CHEMICAL",1.25), ("DOSAGE",4.45), ("UNIT",5.55), ("PRICE",6.45), ("COST IMPACT",7.55)]
    for h,x in headers: v54_text(s,h,x,1.32,1.3,0.11,7.2,True,v54_color("accent"))
    colors=[(0,153,80),(245,158,11),(0,92,184),(126,87,194),(0,51,102),(236,72,153),(20,184,166)]
    yy=1.68
    for i,r in enumerate(rows):
        v54_rect(s,0.55,yy-0.07,11.85,0.46,fill=(255,255,255),line=(235,240,246),radius=True)
        col=colors[i%len(colors)]
        v54_rect(s,0.66,yy,0.25,0.25,fill=col,line=col,radius=True)
        v54_text(s,r.get("label",chr(65+i)),0.70,yy+0.065,0.17,0.07,5.8,True,(255,255,255),PP_ALIGN.CENTER)
        v54_text(s,r.get("chemical") or r.get("supplier") or "Chemical",1.25,yy+0.06,2.6,0.12,7.9,True,v54_color("text"))
        v54_text(s,v54_fmt(r.get("amount",0),2),4.45,yy+0.06,0.85,0.12,7.7,False,v54_color("text"))
        v54_text(s,r.get("unit","g/l"),5.55,yy+0.06,0.65,0.12,7.7,False,v54_color("text"))
        v54_text(s,v54_fmt(r.get("price",0),2),6.45,yy+0.06,0.8,0.12,7.7,False,v54_color("text"))
        v54_text(s,"€0.00",7.55,yy+0.06,0.8,0.12,7.9,True,v54_color("navy"))
        yy += 0.50
    if not rows:
        v54_text(s,"No chemical dosing was defined for this project.",0.7,2.2,6,0.18,12,False,v54_color("muted"))
    v54_rect(s,0.35,6.17,12.45,0.62,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,"⚗  TOTAL CHEMICAL COST",0.65,6.42,2.3,0.11,8,True,v54_color("navy"))
    v54_text(s,f"€{v54_fmt(d.get('Chemical Cost / batch',0),2)} / batch",4.15,6.40,2.2,0.14,12,True,v54_color("navy"))
    v54_text(s,f"{v54_fmt(d.get('Chemical Cost / kg',0),3)} €/kg",9.95,6.40,2.2,0.14,12,True,v54_color("navy"))
    v54_footer(s,3)
    # Slide 4 cost dashboard
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,4,"Cost & Utilities Dashboard","Batch & Per Kg Analysis",project)
    v54_rect(s,0.35,1.05,12.45,0.62,fill=v54_color("navy"),line=v54_color("navy"),radius=True)
    v54_text(s,f"TOTAL COST / BATCH  €{v54_fmt(d.get('Total Cost / batch',0),1)}",0.62,1.27,2.7,0.13,8.5,True,(255,255,255))
    v54_text(s,f"TOTAL COST / KG  €{v54_fmt(d.get('Total Cost / kg',0),3)}",3.55,1.27,2.7,0.13,8.5,True,(255,255,255))
    v54_text(s,f"ENERGY / KG  {v54_fmt(d.get('Energy kWh / kg',0),2)} kWh",6.55,1.27,2.7,0.13,8.5,True,(255,255,255))
    v54_text(s,f"WATER / KG  {v54_fmt(d.get('Water L / kg',0),1)} L",9.55,1.27,2.7,0.13,8.5,True,(255,255,255))
    comps=v54_cost_components(d); maxc=max([float(v or 0) for _,v,_ in comps] or [1])
    v54_rect(s,0.35,1.92,5.80,3.15,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,"Batch Basis (Per Batch)",0.55,2.13,2.3,0.13,9,True,v54_color("navy"))
    yy=2.45
    for lab,val,col in comps:
        v54_progress(s,lab,float(val or 0),maxc,0.65,yy,5.05,col,f"€{v54_fmt(val,2)}")
        yy+=0.36
    v54_text(s,f"TOTAL    €{v54_fmt(d.get('Total Cost / batch',0),1)}",0.65,4.67,2.8,0.14,8.5,True,v54_color("navy"))
    v54_rect(s,6.35,1.92,2.6,3.15,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_donut(s,comps,6.75,2.18,1.0)
    v54_rect(s,9.18,1.92,3.62,3.15,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,"Key Takeaway",9.45,2.15,2.2,0.13,9,True,v54_color("navy"))
    major=max(comps,key=lambda x:float(x[1] or 0))
    v54_text(s,f"{major[0]} is the major cost driver. Optimizing this component will deliver the highest savings potential.",9.45,2.55,2.9,0.85,10,False,v54_color("text"))
    v54_rect(s,0.35,5.35,12.45,0.95,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,"Per Kg Basis",0.60,5.58,1.6,0.13,9,True,v54_color("navy"))
    x=2.05
    for lab,val,col in comps[:6]:
        kg=float(val or 0)/fabric
        v54_text(s,f"{lab}\n€{v54_fmt(kg,3)}",x,5.60,1.50,0.35,7,True,v54_color("text"),PP_ALIGN.CENTER)
        x+=1.72
    v54_footer(s,4)
    # Slide 5 sustainability
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,5,"Sustainability & Carbon Footprint","Environmental Impact Overview",project)
    v54_kpi(s,"Electricity CO₂",v54_fmt(d.get("Electricity CO₂ (kg/batch)",0),1),"kg / batch",0.55,1.20,2.6,1.25,"⚡")
    v54_kpi(s,"Heating CO₂",v54_fmt(d.get("Heating CO₂ (kg/batch)",0),1),"kg / batch",3.45,1.20,2.6,1.25,"♨",True)
    v54_kpi(s,"Total CO₂",v54_fmt(d.get("Total CO₂ (kg/batch)",0),1),f"kg / batch | {v54_fmt(d.get('Total CO₂ / kg (g)',0),1)} g/kg",6.35,1.20,2.6,1.25,"☁")
    v54_kpi(s,"Water",v54_fmt(d.get("Total Water L / batch",0),0),f"L / batch | {v54_fmt(d.get('Water L / kg',0),1)} L/kg",9.25,1.20,2.6,1.25,"◊")
    v54_rect(s,0.55,3.0,12.1,2.1,fill=(255,255,255),line=v54_color("line"),radius=True)
    total_co2=max(float(d.get("Total CO₂ (kg/batch)",0) or 0),0.001)
    e=float(d.get("Electricity CO₂ (kg/batch)",0) or 0); h=float(d.get("Heating CO₂ (kg/batch)",0) or 0)
    v54_text(s,"Impact Distribution (Per Batch)",0.78,3.24,2.5,0.13,9,True,v54_color("navy"))
    v54_progress(s,"Electricity CO₂",e,total_co2,0.82,3.72,10.4,(0,150,90),f"{100*e/total_co2:.1f}%")
    v54_progress(s,"Heating CO₂",h,total_co2,0.82,4.22,10.4,v54_color("accent"),f"{100*h/total_co2:.1f}%")
    v54_rect(s,0.55,5.45,12.1,0.82,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(s,"Key Insight",0.80,5.68,1.4,0.13,9,True,(0,130,80))
    v54_text(s,f"This process generates {v54_fmt(d.get('Total CO₂ / kg (g)',0),1)} g CO₂ per kg and consumes {v54_fmt(d.get('Energy kWh / kg',0),2)} kWh energy and {v54_fmt(d.get('Water L / kg',0),1)} L water per kg.",2.10,5.67,9.8,0.13,8.3,False,v54_color("text"))
    v54_footer(s,5)
    prs.save(out)
    try: graph_path.unlink()
    except Exception: pass

def v54_compare_result(a,b,tol=0.0005):
    # v58_COMPARE_LABOUR_FIX: tolerance reduced so close but visibly different
    # per-kg costs such as Labour/kg 0.023 vs 0.019 are no longer shown as Equal.
    try:
        a=float(a); b=float(b)
        if abs(a-b)<=tol: return "Equal", "0.0%", (110,120,130)
        pct=((b-a)/a*100) if abs(a)>tol else (100.0 if b>a else -100.0)
        better="P1" if a<b else "P2"
        return f"Better: {better}", f"{pct:+.1f}%", (0,150,90) if better=="P1" else v54_color("accent")
    except Exception:
        return "-","-",(110,120,130)

def v54_compare_card(slide,title,p1,p2,x,y,w,h):
    status,pct,col=v54_compare_result(p1,p2)
    v54_rect(slide,x,y,w,h,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_text(slide,title.upper(),x+0.14,y+0.13,w-0.28,0.10,6.7,True,v54_color("muted"))
    v54_text(slide,f"P1  {v54_fmt(p1,3)}",x+0.14,y+0.42,w*0.42,0.16,10.2,True,v54_color("navy"))
    v54_text(slide,f"P2  {v54_fmt(p2,3)}",x+w*0.52,y+0.42,w*0.42,0.16,10.2,True,v54_color("navy"))
    v54_text(slide,pct,x+0.14,y+0.74,w*0.42,0.13,9,True,col)
    v54_text(slide,status,x+w*0.52,y+0.74,w*0.42,0.13,8.2,True,col,PP_ALIGN.RIGHT)

def v54_build_compare_ppt(projects, out, chart_svgs=None, chart_pngs=None):
    p1,p2=projects[0],projects[1]
    r1,r2=v46_calc(p1),v46_calc(p2); c1,c2=r1.get("dashboard",{}),r2.get("dashboard",{})
    chart_pngs=chart_pngs or []
    g1=v46_png_dataurl_to_file(chart_pngs[0],"v54_cmp1") if len(chart_pngs)>0 else None
    g2=v46_png_dataurl_to_file(chart_pngs[1],"v54_cmp2") if len(chart_pngs)>1 else None
    if not g1:
        g1=V38_GENERATED / f"v54_cmp1_{int(time.time()*1000)}.png"
        try:
            v46_render_process_graph(p1,g1,title=p1.get("project_name","Project 1"),width=6.4,height=4.0)
        except Exception:
            create_chart_png(p1,g1)
    if not g2:
        g2=V38_GENERATED / f"v54_cmp2_{int(time.time()*1000)}.png"
        try:
            v46_render_process_graph(p2,g2,title=p2.get("project_name","Project 2"),width=6.4,height=4.0)
        except Exception:
            create_chart_png(p2,g2)
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    # Slide 1 dual executive
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,1,"Executive Compare","Project input and KPI comparison",p1)
    for idx,(p,c,x0) in enumerate([(p1,c1,0.55),(p2,c2,7.0)]):
        m=p.get("machine",{}) or {}; u=p.get("utilities",{}) or {}
        v54_rect(s,x0,1.05,5.75,4.85,fill=(255,255,255),line=v54_color("line"),radius=True)
        v54_text(s,p.get("project_name",f"Project {idx+1}"),x0+0.2,1.25,3.4,0.22,15,True,v54_color("navy"))
        v54_text(s,f"{p.get('company_name','')} | {p.get('process_type','')}",x0+0.22,1.58,3.8,0.12,7,False,v54_color("muted"))
        v54_kpi(s,"Cost/kg",f"€{v54_fmt(c.get('Total Cost / kg',0),3)}","",x0+0.22,2.0,1.58,1.15,"●",idx==0)
        v54_kpi(s,"CO₂/kg",v54_fmt(c.get('Total CO₂ / kg (g)',0),1),"g/kg",x0+1.98,2.0,1.58,1.15,"☁")
        v54_kpi(s,"Water/kg",v54_fmt(c.get('Water L / kg',0),1),"L/kg",x0+3.74,2.0,1.58,1.15,"◊")
        rows=[("Fabric",p.get('fabric_kg',0)),("Flote",p.get('flote',0)),("Machine",m.get('machine_name','-')),("Heating",u.get('heating_source','-'))]
        yy=3.55
        for k,v in rows:
            v54_text(s,k,x0+0.28,yy,1.0,0.11,7,False,v54_color("muted")); v54_text(s,v,x0+1.25,yy,1.7,0.11,7,True,v54_color("text")); yy+=0.32
    v54_text(s,"VS",6.33,3.27,0.45,0.24,18,True,v54_color("accent"),PP_ALIGN.CENTER)
    v54_footer(s,1)
    # Slide 2 process graph compare
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,2,"Process Graph Compare","Two process graphs side by side with chemical legends",p1)
    v54_rect(s,0.35,1.05,6.25,4.35,fill=(255,255,255),line=v54_color("line"),radius=True)
    v54_rect(s,6.75,1.05,6.25,4.35,fill=(255,255,255),line=v54_color("line"),radius=True)
    s.shapes.add_picture(str(g1),Inches(0.55),Inches(1.33),width=Inches(5.85),height=Inches(3.6))
    s.shapes.add_picture(str(g2),Inches(6.95),Inches(1.33),width=Inches(5.85),height=Inches(3.6))
    v54_legend_panel(s,r1.get("chemical_legend",[]),0.55,5.62,5.85,0.92,"CHEMICALS – PROJECT 1")
    v54_legend_panel(s,r2.get("chemical_legend",[]),6.95,5.62,5.85,0.92,"CHEMICALS – PROJECT 2")
    v54_footer(s,2)
    # Slide 3 delta dashboard
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,3,"Cost Delta Dashboard","Batch and kg-based cost comparison",p1)
    metrics=[("Total Cost/kg",c1.get('Total Cost / kg',0),c2.get('Total Cost / kg',0)),("Heating/kg",c1.get('Heating Cost / kg',0),c2.get('Heating Cost / kg',0)),("Electricity/kg",c1.get('Electricity Cost / kg',0),c2.get('Electricity Cost / kg',0)),("Water/kg",c1.get('Water Cost / kg',0),c2.get('Water Cost / kg',0)),("Labour/kg",c1.get('Labour Cost / kg',0),c2.get('Labour Cost / kg',0)),("Chemical/kg",c1.get('Chemical Cost / kg',0),c2.get('Chemical Cost / kg',0))]
    for i,(lab,a,b) in enumerate(metrics):
        v54_compare_card(s,lab,a,b,0.55+(i%3)*4.18,1.20+(i//3)*1.55,3.75,1.18)
    v54_rect(s,0.55,4.80,12.20,1.15,fill=(255,255,255),line=v54_color("line"),radius=True)
    best="Project 1" if c1.get('Total Cost / kg',0)<=c2.get('Total Cost / kg',0) else "Project 2"
    v54_text(s,"Executive Decision",0.85,5.08,2.1,0.13,9,True,v54_color("navy"))
    v54_text(s,f"{best} shows the better total cost per kg in this comparison. Use the delta cards above to identify which utility or cost driver creates the difference.",2.75,5.05,8.9,0.16,9,False,v54_color("text"))
    v54_footer(s,3)
    # Slide 4 carbon compare
    s=prs.slides.add_slide(blank); v54_bg(s); v54_header(s,4,"Carbon Intelligence Compare","Batch and kg-based environmental comparison",p1)
    mets=[("Total CO₂/kg",c1.get('Total CO₂ / kg (g)',0),c2.get('Total CO₂ / kg (g)',0)),("Electricity CO₂/kg",c1.get('Electricity CO₂ / kg (g)',0),c2.get('Electricity CO₂ / kg (g)',0)),("Heating CO₂/kg",c1.get('Heating CO₂ / kg (g)',0),c2.get('Heating CO₂ / kg (g)',0)),("Energy/kg",c1.get('Energy kWh / kg',0),c2.get('Energy kWh / kg',0)),("Water/kg",c1.get('Water L / kg',0),c2.get('Water L / kg',0)),("Time",c1.get('Total Time (min)',0),c2.get('Total Time (min)',0))]
    for i,(lab,a,b) in enumerate(mets):
        v54_compare_card(s,lab,a,b,0.55+(i%3)*4.18,1.20+(i//3)*1.55,3.75,1.18)
    v54_rect(s,0.55,4.80,12.20,1.15,fill=(255,255,255),line=v54_color("line"),radius=True)
    best="Project 1" if c1.get('Total CO₂ / kg (g)',0)<=c2.get('Total CO₂ / kg (g)',0) else "Project 2"
    v54_text(s,"Sustainability Insight",0.85,5.08,2.2,0.13,9,True,(0,130,80))
    v54_text(s,f"{best} shows the better CO₂ intensity per kg. The comparison separates electricity, heating, energy and water impacts for management review.",3.00,5.05,8.75,0.16,9,False,v54_color("text"))
    v54_footer(s,4)
    prs.save(out)
    for g in [g1,g2]:
        try: g.unlink()
        except Exception: pass

def v46_png_dataurl_to_file(data_url, prefix="chart"):
    """Save browser-rendered PNG chart data URL to file for PowerPoint.
    This avoids server-side SVG rendering dependencies and keeps PPT charts identical to the program view.
    """
    if not data_url or not isinstance(data_url, str) or "base64," not in data_url:
        return None
    V38_GENERATED.mkdir(exist_ok=True)
    out = V38_GENERATED / f"{prefix}_{int(time.time()*1000)}.png"
    try:
        raw = data_url.split("base64,", 1)[1]
        out.write_bytes(base64.b64decode(raw))
        return out
    except Exception:
        return None

def v46_svg_to_png(svg_text, prefix="chart"):
    """Legacy fallback disabled when CairoSVG is not installed.
    Use chart_png from the browser for identical chart output.
    """
    return None

def v46_build_single_ppt(project, out, chart_svg=None, chart_png=None):
    data = v46_calc(project)
    d = data.get("dashboard", {})
    # Prefer the browser/program SVG chart so PPT uses the exact same graph style.
    graph_path = v46_png_dataurl_to_file(chart_png, "graph_single_ui") if chart_png else None
    if not graph_path and chart_svg:
        graph_path = v46_svg_to_png(chart_svg, "graph_single_ui")
    if graph_path:
        gdata = v46_calc(project)
    else:
        graph_path = V38_GENERATED / f"graph_single_{int(time.time()*1000)}.png"
        gdata = v46_render_process_graph(project, graph_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 Basic + Machine
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Basic Inputs + Machine", "Executive project and machine overview")
    v46_add_text(s, project.get("project_name","Project"), 0.55, 0.95, 6.0, 0.28, 17, True, (8,55,99))
    v46_add_text(s, f"Company: {project.get('company_name','-')}   |   Process: {project.get('process_type','-')}", 0.55, 1.28, 8.0, 0.20, 9, False, (80,95,115))
    m = project.get("machine", {}) or {}
    u = project.get("utilities", {}) or {}
    total_power = v46_num(m.get("circulation_pump_power"))*v46_num(m.get("pump_ratio"),1)+v46_num(m.get("number_of_reel"))*v46_num(m.get("reel_power"))+v46_num(m.get("fan_power"))*v46_num(m.get("fan_ratio"),1)
    basics = [("Fabric kg", project.get("fabric_kg",0)), ("Flote", project.get("flote",0)), ("Heating", u.get("heating_source","-")), ("Currency", project.get("cost_currency","EUR")), ("Carry Over", project.get("carry_over",0)), ("Fabric Status", project.get("fabric_status","-"))]
    machines = [("Machine", m.get("machine_name","-")), ("Capacity kg", m.get("capacity_kg",0)), ("Drain min", m.get("drain_time_min",0)), ("Pump kW", m.get("circulation_pump_power",0)), ("Reel Power", m.get("reel_power",0)), ("Total Power", round(total_power,2))]
    v46_add_text(s, "Basic Inputs", 0.65, 1.78, 3, 0.22, 13, True, (8,55,99))
    v46_add_text(s, "Machine Parameters", 6.95, 1.78, 3, 0.22, 13, True, (8,55,99))
    for i,(k,v) in enumerate(basics):
        v46_card(s, k, str(v), 0.65+(i%2)*2.95, 2.15+(i//2)*0.92, 2.62, 0.72)
    for i,(k,v) in enumerate(machines):
        v46_card(s, k, str(v), 6.95+(i%2)*2.95, 2.15+(i//2)*0.92, 2.62, 0.72)
    v46_card(s, "Total Cost / kg", str(d.get("Total Cost / kg",0)), 0.65, 5.55, 2.62, 0.72, True)
    v46_card(s, "CO₂ / kg", f'{d.get("Total CO₂ / kg (g)",0)} g', 3.6, 5.55, 2.62, 0.72, True)
    v46_card(s, "Total Time", f'{d.get("Total Time (min)",0)} min', 6.55, 5.55, 2.62, 0.72, True)
    v46_card(s, "Water L/kg", str(d.get("Water L / kg",0)), 9.5, 5.55, 2.62, 0.72, True)

    # Slide 2 graph + chemicals
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Process Graph + Chemicals", "Same graph engine as program view")
    s.shapes.add_picture(str(graph_path), Inches(0.45), Inches(0.95), width=Inches(12.4), height=Inches(4.55))
    legend = "\n".join(gdata.get("chemical_legend", [])[:24]) or "No chemicals"
    v46_rect(s, 0.55, 5.62, 12.25, 1.18, fill=(255,255,255), line=(214,226,238), radius=True)
    v46_add_text(s, "Chemical Legend", 0.75, 5.82, 2.0, 0.18, 10, True, (8,55,99))
    v46_add_text(s, legend, 2.20, 5.76, 10.35, 0.90, 7.4, False, (42,55,70))

    # Slide 3 cost dashboard
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Batch + Kg Cost Dashboard", "Heating, electricity, chemical, water and labour")
    batch = [("Heating / batch", "Heating Cost / batch"), ("Electricity / batch", "Electricity Cost / batch"), ("Chemical / batch", "Chemical Cost / batch"), ("Water / batch", "Water Cost / batch"), ("Waste / batch", "Waste Water Cost / batch"), ("Labour / batch", "Labour Cost / batch"), ("Total / batch", "Total Cost / batch")]
    perkg = [("Heating / kg", "Heating Cost / kg"), ("Electricity / kg", "Electricity Cost / kg"), ("Chemical / kg", "Chemical Cost / kg"), ("Water / kg", "Water Cost / kg"), ("Waste / kg", "Waste Water Cost / kg"), ("Labour / kg", "Labour Cost / kg"), ("Total / kg", "Total Cost / kg")]
    v46_add_text(s, "Batch basis", 0.65, 1.05, 3, 0.2, 13, True, (8,55,99))
    v46_add_text(s, "Kg basis", 6.95, 1.05, 3, 0.2, 13, True, (8,55,99))
    for i,(title,key) in enumerate(batch):
        v46_card(s, title, str(d.get(key,0)), 0.65+(i%2)*2.95, 1.48+(i//2)*0.86, 2.62, 0.68, i==6)
    for i,(title,key) in enumerate(perkg):
        v46_card(s, title, str(d.get(key,0)), 6.95+(i%2)*2.95, 1.48+(i//2)*0.86, 2.62, 0.68, i==6)

    # Slide 4 carbon
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Carbon Footprint Dashboard", "Batch and kg-based footprint indicators")
    carbon_cards = [
        ("Electricity CO₂ / batch", f'{d.get("Electricity CO₂ (kg/batch)",0)} kg'),
        ("Heating CO₂ / batch", f'{d.get("Heating CO₂ (kg/batch)",0)} kg'),
        ("Total CO₂ / batch", f'{d.get("Total CO₂ (kg/batch)",0)} kg'),
        ("Electricity CO₂ / kg", f'{d.get("Electricity CO₂ / kg (g)",0)} g'),
        ("Heating CO₂ / kg", f'{d.get("Heating CO₂ / kg (g)",0)} g'),
        ("Total CO₂ / kg", f'{d.get("Total CO₂ / kg (g)",0)} g'),
        ("Energy / kg", f'{d.get("Energy kWh / kg",0)} kWh'),
        ("Water / kg", f'{d.get("Water L / kg",0)} L'),
    ]
    for i,(k,v) in enumerate(carbon_cards):
        v46_card(s, k, v, 0.75+(i%4)*3.05, 1.25+(i//4)*1.25, 2.65, 0.92, i in [2,5])
    v46_rect(s, 0.75, 4.25, 11.85, 1.25, fill=(255,255,255), line=(214,226,238), radius=True)
    v46_add_text(s, "Interpretation", 0.95, 4.48, 2.0, 0.18, 12, True, (8,55,99))
    v46_add_text(s, "This slide summarizes batch-level and kg-level carbon intensity for management comparison and process improvement decisions.", 0.95, 4.88, 11.3, 0.30, 13, False, (42,55,70))
    prs.save(out)
    try: graph_path.unlink()
    except Exception: pass

def v46_compare_value(d, key):
    return d.get(key, 0)

def v46_build_compare_ppt(projects, out, chart_svgs=None, chart_pngs=None):
    p1, p2 = projects[0], projects[1]
    c1, c2 = v46_calc(p1)["dashboard"], v46_calc(p2)["dashboard"]
    chart_svgs = chart_svgs or []
    chart_pngs = chart_pngs or []
    graph1 = v46_png_dataurl_to_file(chart_pngs[0], "compare_graph1_ui") if len(chart_pngs) > 0 else None
    graph2 = v46_png_dataurl_to_file(chart_pngs[1], "compare_graph2_ui") if len(chart_pngs) > 1 else None
    if not graph1 and len(chart_svgs) > 0:
        graph1 = v46_svg_to_png(chart_svgs[0], "compare_graph1_ui")
    if not graph2 and len(chart_svgs) > 1:
        graph2 = v46_svg_to_png(chart_svgs[1], "compare_graph2_ui")
    if graph1:
        g1 = v46_calc(p1)
    else:
        graph1 = V38_GENERATED / f"compare_graph1_{int(time.time()*1000)}.png"
        g1 = v46_render_process_graph(p1, graph1, title=p1.get("project_name","Project 1"), width=6.4, height=4.0)
    if graph2:
        g2 = v46_calc(p2)
    else:
        graph2 = V38_GENERATED / f"compare_graph2_{int(time.time()*1000)}.png"
        g2 = v46_render_process_graph(p2, graph2, title=p2.get("project_name","Project 2"), width=6.4, height=4.0)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def compare_result(a, b, tol=0.005):
        try:
            a=float(a); b=float(b)
            if abs(a-b) <= tol:
                return "0.0% | Equal", (110, 120, 130)
            d=((b-a)/a*100) if abs(a) > tol else (100.0 if b>a else -100.0)
            better = "P1" if a < b else "P2"
            # Lower is better for cost, energy, water, time and carbon metrics.
            return f"{d:+.1f}% | Better: {better}", (0, 140, 72)
        except Exception:
            return "-", (110, 120, 130)

    def diff_text(a, b):
        return compare_result(a, b)[0]

    def diff_color(a, b):
        return compare_result(a, b)[1]

    # Slide 1 Basic + Machine compare
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Basic Inputs + Machine Compare", "Project input and machine parameter comparison")
    for idx, p in enumerate([p1,p2]):
        x0 = 0.65 if idx==0 else 6.95
        m = p.get("machine",{}) or {}
        u = p.get("utilities",{}) or {}
        v46_add_text(s, p.get("project_name", f"Project {idx+1}"), x0, 1.0, 4.5, 0.25, 16, True, (8,55,99))
        basics = [("Fabric kg", p.get("fabric_kg",0)), ("Flote", p.get("flote",0)), ("Process", p.get("process_type","-")), ("Heating", u.get("heating_source","-"))]
        machines = [("Machine", m.get("machine_name","-")), ("Capacity", m.get("capacity_kg",0)), ("Drain min", m.get("drain_time_min",0)), ("Pump kW", m.get("circulation_pump_power",0))]
        for i,(k,v) in enumerate(basics+machines):
            v46_card(s, k, str(v), x0+(i%2)*2.75, 1.55+(i//2)*0.82, 2.45, 0.62)
    # Clean compare cover: no repeated KPI summary boxes at the bottom.
    v46_rect(s, 5.95, 1.55, 0.72, 4.20, fill=(255,247,239), line=(255,215,176), radius=True)
    v46_add_text(s, "VS", 6.07, 3.45, 0.46, 0.24, 16, True, (255,122,0), PP_ALIGN.CENTER)

    # Slide 2 dual graph + chemicals
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Process Graph Compare", "Two process graphs side by side with chemical legends")
    s.shapes.add_picture(str(graph1), Inches(0.55), Inches(1.0), width=Inches(5.95), height=Inches(3.65))
    s.shapes.add_picture(str(graph2), Inches(6.85), Inches(1.0), width=Inches(5.95), height=Inches(3.65))
    v46_rect(s, 0.55, 4.95, 5.95, 1.35, fill=(255,255,255), line=(214,226,238), radius=True)
    v46_rect(s, 6.85, 4.95, 5.95, 1.35, fill=(255,255,255), line=(214,226,238), radius=True)
    v46_add_text(s, "Chemicals – Project 1", 0.75, 5.12, 2.3, 0.18, 9, True, (8,55,99))
    v46_add_text(s, "\n".join(g1.get("chemical_legend", [])[:10]) or "No chemicals", 0.75, 5.38, 5.5, 0.7, 6.8, False, (42,55,70))
    v46_add_text(s, "Chemicals – Project 2", 7.05, 5.12, 2.3, 0.18, 9, True, (8,55,99))
    v46_add_text(s, "\n".join(g2.get("chemical_legend", [])[:10]) or "No chemicals", 7.05, 5.38, 5.5, 0.7, 6.8, False, (42,55,70))

    # Slide 3 cost compare
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Batch + Kg Cost Compare", "Heating, electricity, chemical, water and labour comparison")
    metrics = [
        ("Heating", "Heating Cost / batch", "Heating Cost / kg"),
        ("Electricity", "Electricity Cost / batch", "Electricity Cost / kg"),
        ("Chemical", "Chemical Cost / batch", "Chemical Cost / kg"),
        ("Water", "Water Cost / batch", "Water Cost / kg"),
        ("Labour", "Labour Cost / batch", "Labour Cost / kg"),
        ("Total", "Total Cost / batch", "Total Cost / kg"),
    ]
    v46_add_text(s, "Metric", 0.75, 1.05, 1.4, 0.2, 10, True, (8,55,99))
    v46_add_text(s, "Project 1 batch / kg", 2.55, 1.05, 2.6, 0.2, 10, True, (8,55,99))
    v46_add_text(s, "Project 2 batch / kg", 5.75, 1.05, 2.6, 0.2, 10, True, (8,55,99))
    v46_add_text(s, "Difference", 9.05, 1.05, 2.2, 0.2, 10, True, (8,55,99))
    for i,(lab,bk,kgk) in enumerate(metrics):
        y=1.45+i*0.73
        v46_rect(s,0.65,y,12.05,0.55,fill=(255,255,255),line=(225,232,240),radius=True)
        v46_add_text(s, lab, 0.85, y+0.17, 1.4, 0.14, 9, True, (55,70,90))
        v46_add_text(s, f'{c1.get(bk,0)} / {c1.get(kgk,0)}', 2.55, y+0.17, 2.6, 0.14, 9, False, (55,70,90), PP_ALIGN.CENTER)
        v46_add_text(s, f'{c2.get(bk,0)} / {c2.get(kgk,0)}', 5.75, y+0.17, 2.6, 0.14, 9, False, (55,70,90), PP_ALIGN.CENTER)
        v46_add_text(s, diff_text(c1.get(kgk,0), c2.get(kgk,0)), 9.05, y+0.17, 2.8, 0.14, 8, True, diff_color(c1.get(kgk,0), c2.get(kgk,0)), PP_ALIGN.CENTER)

    # Slide 4 carbon compare
    s = prs.slides.add_slide(blank); v46_bg(s); v46_header(s, "Carbon Footprint Compare", "Batch and kg-based carbon comparison")
    carbon = [
        ("Electricity CO₂", "Electricity CO₂ (kg/batch)", "Electricity CO₂ / kg (g)", "kg CO₂/batch", "g CO₂/kg"),
        ("Heating CO₂", "Heating CO₂ (kg/batch)", "Heating CO₂ / kg (g)", "kg CO₂/batch", "g CO₂/kg"),
        ("Total CO₂", "Total CO₂ (kg/batch)", "Total CO₂ / kg (g)", "kg CO₂/batch", "g CO₂/kg"),
        ("Energy", "Electricity (kWh/batch)", "Energy kWh / kg", "kWh/batch", "kWh/kg"),
        ("Water", "Total Water L / batch", "Water L / kg", "L/batch", "L/kg"),
    ]
    for i,(lab,bk,kgk,bunit,kgunit) in enumerate(carbon):
        y=1.35+i*0.85
        v46_rect(s,0.75,y,11.8,0.62,fill=(255,255,255),line=(225,232,240),radius=True)
        v46_add_text(s, lab, 0.95, y+0.20, 1.5, 0.15, 10, True, (8,55,99))
        v46_add_text(s, f'P1: {c1.get(bk,0)} {bunit} | {c1.get(kgk,0)} {kgunit}', 2.55, y+0.20, 3.3, 0.15, 8.5, False, (55,70,90))
        v46_add_text(s, f'P2: {c2.get(bk,0)} {bunit} | {c2.get(kgk,0)} {kgunit}', 5.95, y+0.20, 3.3, 0.15, 8.5, False, (55,70,90))
        v46_add_text(s, diff_text(c1.get(kgk,0), c2.get(kgk,0)), 9.65, y+0.20, 2.25, 0.15, 8, True, diff_color(c1.get(kgk,0), c2.get(kgk,0)), PP_ALIGN.CENTER)

    prs.save(out)
    for p in [graph1, graph2]:
        try: p.unlink()
        except Exception: pass

@app.post("/api/export/chart-png")
async def export_chart(req: Request):
    try:
        payload = await req.json()
        chart_png = payload.get("chart_png") if isinstance(payload, dict) else None
        project = payload.get("project", payload) if isinstance(payload, dict) else payload
        out = V38_GENERATED / f"DyeFlow_RS_Chart_{int(time.time())}.png"
        if chart_png and isinstance(chart_png, str) and chart_png.startswith("data:image/png;base64,"):
            out.write_bytes(base64.b64decode(chart_png.split(",",1)[1]))
        else:
            v46_render_process_graph(project, out)
        return FileResponse(out, media_type="image/png", filename="DyeFlow_RS_Chart.png")
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/export/powerpoint")
async def export_ppt(req: Request):
    try:
        payload = await req.json()
        project = payload.get("project", payload) if isinstance(payload, dict) else payload
        chart_svg = payload.get("chart_svg") if isinstance(payload, dict) else None
        chart_png = payload.get("chart_png") if isinstance(payload, dict) else None
        V38_GENERATED.mkdir(exist_ok=True)
        out = V38_GENERATED / f"DyeFlow_RS_Report_{int(time.time())}.pptx"
        v54_build_single_ppt(project, out, chart_svg=chart_svg, chart_png=chart_png)
        return FileResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="DyeFlow_RS_Report.pptx")
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/export/comparison-ppt")
async def comparison_ppt(req: Request):
    try:
        payload = await req.json()
        projects = payload.get("projects") or payload.get("selectedProjects") or []
        if not projects and payload.get("project1") and payload.get("project2"):
            projects = [payload.get("project1"), payload.get("project2")]
        if len(projects) < 2:
            raise HTTPException(400, "Two projects required for compare report")
        chart_svgs = payload.get("chart_svgs") or []
        chart_pngs = payload.get("chart_pngs") or []
        out = V38_GENERATED / f"DyeFlow_RS_Compare_{int(time.time())}.pptx"
        v54_build_compare_ppt(projects[:2], out, chart_svgs=chart_svgs[:2], chart_pngs=chart_pngs[:2])
        return FileResponse(out, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="DyeFlow_RS_Compare_Report.pptx")
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})
        p.mkdir(exist_ok=True)
        
